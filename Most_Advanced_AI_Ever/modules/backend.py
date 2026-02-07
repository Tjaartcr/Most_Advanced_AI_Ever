import ollama
import re
import os
import json
import time
import base64
import mimetypes
import tempfile
from datetime import datetime
from werkzeug.utils import secure_filename

import whisper
from flask import Flask, request, jsonify, send_from_directory
from flask_socketio import SocketIO, emit, join_room, leave_room

# Optional libs (may not be installed in every env)
try:
    import chromadb
except Exception:
    chromadb = None
try:
    import ollama
except Exception:
    ollama = None

from listenWEBUI import WEBUIListenModule
from auth_manager import login as do_login, signup as do_signup, get_users, delete_user
from query_logger import log_user_query, log_user_response, read_logs

import chardet
from pdfminer.high_level import extract_text
from pdf2image import convert_from_path
import pytesseract
import base64
import io
import PyPDF2
from PyPDF2 import PdfFileReader
import pdfplumber
from langchain_chroma import Chroma

import asyncio
import tempfile
from flask import send_file   # modify existing import line to include send_file if needed

import io, re, base64, mimetypes, tempfile, os, time
import chardet
from werkzeug.utils import secure_filename

import io, re, base64, mimetypes, tempfile, os, time
import chardet
from werkzeug.utils import secure_filename
from datetime import datetime
        
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TMP_DIR = os.path.join(BASE_DIR, 'tmp')
os.makedirs(TMP_DIR, exist_ok=True)

# Initialize listener and Whisper
webui_listener = WEBUIListenModule()
whisper_model = whisper.load_model("base.en")

app = Flask(__name__, static_folder='./webui-src/dist', static_url_path='')
socketio = SocketIO(app, cors_allowed_origins='*', async_mode='threading')


chromadb.api.client.SharedSystemClient.telemetry_enabled = False


# Module-level state
state = {
    'logs': [],
    'queries': [],
    'resp': [],
    'settings': {
        'use_whisper': False,
        'use_vosk': True,
        'enter_submits': True
    }
}

# Session tracking
session_users = {}  # Maps SID -> username

my_file_content = ""

collection = None   # Chroma collection shared across events
os.environ["CHROMA_TELEMETRY_ENABLED"] = "false"

##from langchain.text_splitter import RecursiveCharacterTextSplitter
import PyPDF2


model_selected = ""
thinkbot_model = ""
rag_clear = False
detected_model = ""
RAG_model = ""

import io
import os
import tempfile
import traceback
import base64
import re


# ---------- New helper: robust Ollama embedding helper ----------
def get_ollama_embedding(text: str):
    """
    Return a single embedding vector (list of floats) for `text`.
    Tries multiple shapes/keys and falls back to sentence-transformers if Ollama unavailable or fails.
    Returns: list[float] or None
    """
    if text is None:
        return None

    # Try Ollama if available
    if ollama is not None:
        try:
            # preferred API: ollama.embed(model=..., input=...)
            # some users have older/newer clients; try both 'embed' and 'embed*' helpers
            emb_resp = None
            if hasattr(ollama, "embed"):
                emb_resp = ollama.embed(model="mxbai-embed-large", input=text)
            elif hasattr(ollama, "embed_text"):
                emb_resp = ollama.embed_text(model="mxbai-embed-large", input=text)
            elif hasattr(ollama, "embeddings"):
                # older name - but be sure to pass input/prompt
                try:
                    emb_resp = ollama.embeddings(model="mxbai-embed-large", input=text)
                except TypeError:
                    # some server versions want "prompt"
                    try:
                        emb_resp = ollama.embeddings(model="mxbai-embed-large", prompt=text)
                    except Exception:
                        emb_resp = None
            else:
                # final try: call embed() fallback with prompt
                try:
                    emb_resp = ollama.embed(model="mxbai-embed-large", input=text)
                except Exception:
                    emb_resp = None

            if emb_resp:
                # Response shapes vary: try common keys
                if isinstance(emb_resp, dict):
                    if "embeddings" in emb_resp:
                        emb = emb_resp.get("embeddings")
                    elif "embedding" in emb_resp:
                        emb = emb_resp.get("embedding")
                    elif "data" in emb_resp and isinstance(emb_resp["data"], list) and len(emb_resp["data"]) > 0:
                        # sometimes wrapped
                        first = emb_resp["data"][0]
                        emb = first.get("embedding") or first.get("embeddings") or None
                    else:
                        # sometimes top-level is the embedding array
                        emb = emb_resp.get("emb") or emb_resp.get("vector") or None
                else:
                    emb = emb_resp

                # normalize nested lists
                if emb is None:
                    emb = emb_resp
                # If nested like [[...]] pick first inner if present
                if isinstance(emb, list) and len(emb) > 0 and isinstance(emb[0], list):
                    emb = emb[0]

                # ensure floats
                if isinstance(emb, (list, tuple)):
                    try:
                        emb_f = [float(x) for x in emb]
                        return emb_f
                    except Exception:
                        return None

        except Exception as e:
            print(f"[WARN] Ollama embedding attempt failed: {e}")

    # Ollama failed or not available -> fallback to sentence-transformers if installed
    try:
        from sentence_transformers import SentenceTransformer
        st_model = SentenceTransformer('all-MiniLM-L6-v2')
        vect = st_model.encode(text)
        emb = vect.tolist() if hasattr(vect, "tolist") else list(vect)
        return emb
    except Exception as e:
        print(f"[WARN] sentence-transformers fallback failed: {e}")
        return None


# ======= Query helper (can be called from elsewhere in backend) =======
def query_with_retrieval(prompt: str, n_results: int = 1):
    """Retrieve top chunks from Chroma and run generation with Ollama."""
    try:
        if ollama is None or chromadb is None:
            raise RuntimeError("Ollama or Chroma not available")

        # get a robust embedding for the prompt
        q_vector = get_ollama_embedding(prompt)
        if q_vector is None:
            raise RuntimeError("Embedding generation failed")

        client = chromadb.Client()
        collection = client.get_or_create_collection(name="docs")
        # query expects a list of embeddings
        res = collection.query(query_embeddings=[q_vector], n_results=n_results)
        top_doc = ""
        try:
            # many chroma versions return {'documents': [[doc1, doc2,...]], ...}
            docs_field = res.get("documents", [[]])
            if isinstance(docs_field, list) and len(docs_field) > 0:
                # pick first result from first list
                maybe = docs_field[0]
                if isinstance(maybe, list) and len(maybe) > 0:
                    top_doc = maybe[0]
                elif isinstance(maybe, str):
                    top_doc = maybe
        except Exception:
            top_doc = ""

    except Exception as e:
        print(f"[WARN] retrieval failed: {e}")
        top_doc = ""

    gen_prompt = f"Using this data: {top_doc}\n\nRespond to: {prompt}"
    try:
        out = ollama.generate(model="tinyllama", prompt=gen_prompt, stream=False)
        # output shapes vary
        if isinstance(out, dict):
            text = out.get("response") or out.get("text") or out.get("output") or str(out)
        else:
            text = str(out)
    except Exception as e:
        text = f"[generation failed: {e}]"
    return text

def extract_pdf_text_with_ocr(pdf_path):
    try:
        text = extract_text(pdf_path)
        if len(text.strip()) > 10:  # if decent text extracted, return it
            return text
        else:
            # Fallback to OCR if text is too short or empty
            print("[INFO] PDF text extraction empty or too short, using OCR fallback.")
            pages = convert_from_path(pdf_path)
            ocr_text = ""
            for page in pages:
                ocr_text += pytesseract.image_to_string(page)
            return ocr_text
    except Exception as e:
        print(f"[ERROR] PDF extraction + OCR failed: {e}")
        return ""


# Utility: safe save
def _save_tmp(filename, data_bytes):
    safe = secure_filename(filename)
    path = os.path.join(TMP_DIR, safe)
    with open(path, 'wb') as f:
        f.write(data_bytes)
    return path

# Socket handlers
@socketio.on('connect')
def connect():
    emit('full_state', state)

@socketio.on('disconnect')
def disconnect():
    sid = request.sid
    user = session_users.pop(sid, None)
    if user:
        leave_room(user)
        print(f"[DISCONNECT] {sid} left room {user}")


# --- helper: singleton chroma client creator (safe) ---

_CHROMA_CLIENT = None

def get_chroma_client(persist_directory=None, chroma_db_impl="duckdb+parquet"):
    global _CHROMA_CLIENT
    if _CHROMA_CLIENT is not None:
        return _CHROMA_CLIENT
    try:
        import chromadb
        from chromadb.config import Settings as ChromaSettings
    except Exception as e:
        print(f"[RAG] chromadb import failed: {e}")
        return None

    try:
        settings = ChromaSettings(chroma_db_impl=chroma_db_impl,
                                  persist_directory=persist_directory)

        client = None
        try:
##            from chromadb.config import Settings as ChromaSettings
            chroma_settings = ChromaSettings(
                chroma_db_impl="duckdb+parquet",
                persist_directory=os.path.join(BASE_DIR, "chroma_db")
            )
            client = chromadb.Client(chroma_settings)
        except Exception:
            # fallback to default client constructor; may still fail, handle below
            try:
                client = chromadb.Client()
            except Exception as e:
                client = None
                print(f"[RAG] chromadb.Client fallback failed: {e}")


            if client is None:
                client = get_chroma_client(persist_directory=CHROMA_DIR)

        return client




##        _CHROMA_CLIENT = chromadb.Client(settings=settings)
##        print(f"[RAG] Chroma client created with persist_directory={persist_directory}")
##        return _CHROMA_CLIENT
    except Exception as e:
        print(f"[RAG] Embedded Chroma init failed: {e}; trying fallback chromadb.Client()")
        try:
            _CHROMA_CLIENT = chromadb.Client()
            print("[RAG] Chroma client created via fallback chromadb.Client()")
            return _CHROMA_CLIENT
        except Exception as e2:
            print(f"[RAG] chromadb.Client() fallback failed: {e2}")
            _CHROMA_CLIENT = None
            return None

_MODEL_STR_PATTERN = re.compile(r"^[A-Za-z0-9_\-]+(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)?$")

def extract_model(payload) -> str:
    """
    Robustly extract and clean a model name. Always returns a string ("" if none found).

    Rules:
      - If payload is a dict and contains thinkbot_model/model/model_used -> use it.
      - Else if payload is a dict and has 'payload' whose value is a dict -> descend and try again.
      - Else if payload is a plain string and *looks like a model name* (pattern) -> accept it.
      - Otherwise return "".

    Example results:
      {'payload': {'thinkbot_model': 'moondream'}} -> "moondream"
      {'thinkbot_model': 'deepseek-r1-1.5b'}    -> "deepseek-r1"
      "deepseek-r1"                             -> "deepseek-r1"
      {'type': 'log', 'payload': 'I see, a man ...'} -> ""
    """
    model_raw = ""

    # If dict, prefer direct model keys at top-level
    if isinstance(payload, dict):
        # Direct model keys first
        if "detected_model" in payload or "model" in payload or "model_used" in payload:
            model_raw = payload.get("detected_model") or payload.get("model") or payload.get("model_used") or ""
        else:
            # Only descend into 'payload' if it's a dict (avoid descending into plain log strings)
            inner = payload.get("payload")
            if isinstance(inner, dict):
                model_raw = inner.get("detected_model") or inner.get("model") or inner.get("model_used") or ""
            else:
                # No model found in this event
                return ""
    else:
        # payload not a dict; if it's a plain string that *looks like a model* accept it
        if isinstance(payload, str) and _MODEL_STR_PATTERN.match(payload.strip()):
            model_raw = payload.strip()
        else:
            return ""

    model_raw = (model_raw or "").strip()
    if not model_raw:
        return ""

    # Clean suffixes like ":1.5b", "-1.5b", ":8b", "-8b", etc.
    model_clean = re.sub(r"(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)$", "", model_raw)

    return model_clean


@socketio.on('gui_event')
def handle_event(data):
    global client, collection, detected_model, thinkbot_model, model_selected, rag_clear  # Declare globals at the very top
    sid = request.sid
    print(f"[DEBUG] Received gui_event from sid={sid}: {data!r}")

    user = data.get('username') if isinstance(data, dict) else None
    t = data.get('type') if isinstance(data, dict) else None
    payload = data.get('payload') if isinstance(data, dict) else None

    event = data
    print("[EVENT] Received event 'new event' : ", str(event))

    model_used = extract_model(event)
    print(f"[EVENT] Received gui_event 'model_used' : {model_used}")

    if not user:
        user = 'Tjaart'  # fallback
    # Track user per session
    session_users[sid] = user
    join_room(user)
    global my_file_content
    
    # --- LOGIN ---
    if t == 'login':
        print(f"[LOGIN] User: {user}")
        socketio.emit('state_update', {
            'type': 'login_ack',
            'username': user
        }, room=user)
        return

    # --- LOG ---
    if t == 'log':
        # Always produce a structured entry so frontend can read models reliably
        if isinstance(payload, str):
            # Keep old logging side-effect
            log_user_response(user, payload)

            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": "",
                "response": payload,
                "models": {"thinkbot": model_used or ""}
            }

            socketio.emit('state_update', {
                'type': 'log',
                'payload': entry,
                'username': user
            }, room=user)

        elif isinstance(payload, dict):
            # Ensure models is a dict and add thinkbot
            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
            models = dict(models)  # copy-safe
            models["thinkbot"] = model_used or ""
            query = payload.get("query", "")

            if query:
                guery = query.replace(":","")
                print(f"New Query supplied : {query}")
            else:
                print("There is no query supplied...")

            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": guery,
                "response": payload.get("response", ""),
                "models": models
            }

            log_user_query(user, entry)

            socketio.emit('state_update', {
                'type': 'log',
                'payload': entry,
                'username': user
            }, room=user)

            socketio.emit('state_update', {
                'type': 'query',
                'payload': entry,
                'username': user
            }, room=user)

        return

    # --- RESP ---
    if t == 'resp':
        if isinstance(payload, str):
            log_user_response(user, payload)

            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": "",
                "response": payload,
                "models": {"thinkbot": model_used or ""}
            }

            socketio.emit('state_update', {
                'type': 'log',
                'payload': entry,
                'username': user
            }, room=user)

        elif isinstance(payload, dict):
            # Preserve any existing models and inject thinkbot
            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
            models = dict(models)
            models["thinkbot"] = model_used or ""
            query = payload.get("query", "")
            
            if query:
                query = query.replace(":","")
                print(f"New Query supplied : {query}")
            else:
                print("No Query supplied...")
                 
            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": query,
                "response": payload.get("response", ""),
                "models": models
            }

            log_user_query(user, entry)

            socketio.emit('state_update', {
                'type': 'log',
                'payload': entry,
                'username': user
            }, room=user)

            socketio.emit('state_update', {
                'type': 'query',
                'payload': entry,
                'username': user
            }, room=user)

            print(f" \n thinkbot_model : {model_used} \n")

        return

    # --- SETTING ---
    if t == 'setting':
        socketio.emit('state_update', {
            'type': 'setting',
            'payload': payload,
            'username': user
        }, room=user)
        return


    if t in ('setting', 'detected_model'):
        data_payload = data.get('payload')
        user = data.get('username') or 'default_user'
##        new_detected_model = data.get('detected_model')

        # Try to extract the thinkbot model name from different payload shapes:
        model_selected = None
        
        if isinstance(data_payload, dict):
            # common: payload = { 'detected_model': 'tinyllama' }
            if 'detected_model' in data_payload:
                model_selected = data_payload['detected_model']
            else:
                # fallback: payload may be { '<some_model_key>': '<name>' }
                # find any key that ends with '_model' and pick its value
                model_keys = [k for k in data_payload.keys() if k.endswith('_model')]
                if len(model_keys) == 1:
                    model_selected = data_payload[model_keys[0]]
        elif isinstance(data_payload, str):
            # sometimes frontend sends a raw string payload (e.g. from 'think_model' emit)
            model_selected = data_payload

        # If we found a model, assign and notify
        if model_selected:
            # ensure we set the module/global variable (declare `global detected_model` above if needed)
            try:
                global detected_model
            except NameError:
                # if detected_model doesn't exist yet, create it
                detected_model = None

            detected_model = model_selected  # correct assignment (not ==)
            print(f"\n The detected_model for RAG is '{detected_model} \n'")
            print(f"\n The new detected_model SETTINGS is now '{detected_model}' \n")

            # Emit a consistent state_update so clients can sync
            socketio.emit('state_update', {
                'type': 'setting',
                'payload': {'detected_model': detected_model},
                'username': user
            }, room=user)
        else:
            # No model found — still forward the raw payload for other settings
            print(f"[setting] received payload but no model detected: {data_payload}")
            socketio.emit('state_update', {
                'type': 'setting',
                'payload': data_payload,
                'username': user
            }, room=user)

        return


    # --- SETTING ---
    if t == 'think_model':

        payload = data.get('payload')
        user = data.get('username') or 'default_user'
        model_selected = data.get('think_model')

        print(f"The model for RAG 'model_selected' is {model_selected}")
        return


    # --- RESET RAG ---
    if t == 'reset_rag':
        import getpass, traceback, time

        # globals we will clear (try many names used across your app)
        GLOBAL_NAMES_TO_CLEAR = [
            'RAG_STORE', 'rag_store', 'USER_FILES', 'user_files',
            'my_file_content', 'my_file_contents', 'uploaded_doc_ids',
            'rag_enabled', '_CHROMA_CLIENT', '_CHROMA_CLIENT_SETTINGS', 'collection'
        ]

        CHROMA_DIRS_TO_TRY = []
        # primary location (your BASE_DIR)
        CHROMA_DIRS_TO_TRY.append(os.path.join(BASE_DIR, "chroma_db"))
        # alternate locations you used in other codepaths
        CHROMA_DIRS_TO_TRY.append(os.path.join(os.getcwd(), "chroma_db"))
        CHROMA_DIRS_TO_TRY.append(os.path.join(os.getcwd(), "chroma_store"))
        # module-local chroma_db if running from modules folder
        CHROMA_DIRS_TO_TRY.append(os.path.join(os.path.dirname(__file__), "chroma_db"))
        # dedupe
        CHROMA_DIRS_TO_TRY = [os.path.normpath(p) for p in dict.fromkeys([p for p in CHROMA_DIRS_TO_TRY if p])]

        user = data.get('username') if isinstance(data, dict) else user
        user = user or getpass.getuser()
        print(f"[RAG] reset_rag requested by {user}")
        print(f"[RAG] CHROMA_DIRS to try: {CHROMA_DIRS_TO_TRY}")

        # Attempt to import chromadb and silence telemetry to reduce noisy logs
        chromadb_mod = globals().get('chromadb', None)
        try:
            if chromadb_mod is None:
                import chromadb as chromadb_imported
                chromadb_mod = chromadb_imported
                globals()['chromadb'] = chromadb_mod
            try:
                setattr(chromadb_mod, "telemetry_enabled", False)
            except Exception:
                pass
            try:
                telemetry = getattr(chromadb_mod, "telemetry", None)
                if telemetry is not None:
                    for name in ("capture", "capture_event", "send_event", "capture_exception"):
                        if hasattr(telemetry, name):
                            try:
                                setattr(telemetry, name, lambda *a, **kw: None)
                            except Exception:
                                pass
            except Exception:
                pass
        except Exception as e:
            print(f"[RAG] chromadb import/telemetry patch failed: {e}")
            chromadb_mod = None

        # helper: attempt version-tolerant deletion on a client/collection
        def _try_delete_collection_version_tolerant(client, colname="docs"):
            """
            Returns True if we believe persistent deletion succeeded (or store is effectively empty).
            """
            if client is None:
                return False
            deleted = False
            try:
                # preferred high-level delete_collection
                if hasattr(client, "delete_collection"):
                    try:
                        # some versions accept positional, some keyword
                        try:
                            client.delete_collection(colname)
                        except TypeError:
                            client.delete_collection(name=colname)
                        print(f"[RAG] client.delete_collection('{colname}') succeeded")
                        return True
                    except Exception as exc:
                        print(f"[RAG] client.delete_collection failed: {exc!r}")

                # obtain collection object in a version tolerant way
                coll = None
                try:
                    if hasattr(client, "get_or_create_collection"):
                        coll = client.get_or_create_collection(name=colname)
                    else:
                        # try get_collection then create_collection
                        try:
                            if hasattr(client, "get_collection"):
                                coll = client.get_collection(colname)
                        except Exception:
                            coll = None
                        if coll is None and hasattr(client, "create_collection"):
                            coll = client.create_collection(colname)
                except Exception as e:
                    print(f"[RAG] failed to obtain collection object: {e}")
                    coll = None

                if coll is None:
                    # nothing more we can do on this client
                    return False

                # 1) try no-arg delete (supported on some versions)
                try:
                    if hasattr(coll, "delete"):
                        try:
                            coll.delete()
                            print("[RAG] collection.delete() (no-arg) invoked")
                            return True
                        except ValueError as ve:
                            # API demands ids/where; fallback below
                            print(f"[RAG] collection.delete() raised ValueError (expected on some APIs): {ve}")
                        except Exception as exc:
                            print(f"[RAG] collection.delete() (no-arg) failed: {exc!r}")
                except Exception:
                    pass

                # 2) list ids then delete by ids
                try:
                    ids = []
                    try:
                        res = coll.get(include=['ids'])
                        if isinstance(res, dict):
                            ids_field = res.get('ids', [[]])
                        else:
                            ids_field = []
                    except Exception:
                        try:
                            res = coll.get()
                            ids_field = res.get('ids', [[]]) if isinstance(res, dict) else []
                        except Exception:
                            ids_field = []

                    # Normalize shapes like [[id1,id2]] or [id1,id2]
                    if isinstance(ids_field, list) and len(ids_field) > 0 and isinstance(ids_field[0], list):
                        ids = ids_field[0]
                    elif isinstance(ids_field, list):
                        ids = ids_field
                    ids = [str(i) for i in ids if i is not None]

                    if ids:
                        try:
                            coll.delete(ids=ids)
                            print(f"[RAG] collection.delete(ids=...) succeeded ({len(ids)} ids)")
                            return True
                        except Exception as exc:
                            print(f"[RAG] collection.delete(ids=...) failed: {exc!r}")
                    else:
                        # no ids -> treat as empty
                        print("[RAG] collection.get() returned no ids; treating store as empty")
                        return True
                except Exception as exc:
                    print(f"[RAG] listing+delete attempt failed: {exc!r}")

                # 3) try deleting by metadata 'where' shapes (best-effort)
                try:
                    # try a few shapes that different chroma versions expect
                    for where_shape in (
                        {"username": user},
                        {"metadata": {"username": user}},
                        {"metadatas": {"username": user}},
                        {"$or": [{"metadata.username": user}, {"username": user}]},
                    ):
                        try:
                            coll.delete(where=where_shape)
                            print(f"[RAG] collection.delete(where={where_shape}) invoked")
                            return True
                        except Exception:
                            continue
                except Exception:
                    pass

            except Exception as e:
                print(f"[RAG] _try_delete_collection_version_tolerant error: {e}\n{traceback.format_exc()}")

            return deleted

        # helper to create client for a given persist_directory (try to reuse get_chroma_client helper if available)
        def _client_for_path(persist_path):
            # prefer your helper if present
            try:
                helper = globals().get("get_chroma_client")
                if callable(helper):
                    res = helper(persist_directory=persist_path)
                    # helper may return (client, path)
                    if isinstance(res, tuple) and len(res) >= 1:
                        return res[0]
                    return res
            except Exception as e:
                print(f"[RAG] get_chroma_client helper failed for {persist_path}: {e}")

            # fallback direct creation
            try:
                import chromadb
                try:
                    from chromadb.config import Settings as ChromaSettings
                    settings = ChromaSettings(chroma_db_impl="duckdb+parquet", persist_directory=persist_path)
                    client = chromadb.Client(settings)
                    return client
                except Exception:
                    try:
                        client = chromadb.Client()
                        return client
                    except Exception as e2:
                        print(f"[RAG] chromadb.Client() fallback failed for {persist_path}: {e2}")
                        return None
            except Exception as e:
                print(f"[RAG] chromadb import failed in _client_for_path: {e}")
                return None

        persistent_deleted_any = False
        tried_paths = []
        for path in CHROMA_DIRS_TO_TRY:
            try:
                client = _client_for_path(path)
                tried_paths.append((path, bool(client)))
                if client is None:
                    print(f"[RAG] No client for {path}; skipping persistent deletion there.")
                    continue
                ok = _try_delete_collection_version_tolerant(client, colname="docs")
                if ok:
                    print(f"[RAG] Persistent deletion succeeded for path {path}")
                    persistent_deleted_any = True
                else:
                    print(f"[RAG] Persistent deletion did not succeed for path {path} (no error, may be empty or locked).")
            except Exception as e:
                print(f"[RAG] exception trying path {path}: {e}\n{traceback.format_exc()}")

        # Try to clear rag_embeddings module caches if present (best-effort)
        try:
            import rag_embeddings
            FECls = getattr(rag_embeddings, "file_embeddings", None)
            if FECls:
                try:
                    fe_inst = FECls()
                    # try well-named clear methods
                    for name in ("clear_store", "clear", "reset", "build_store"):
                        fn = getattr(fe_inst, name, None)
                        if callable(fn):
                            # if build_store takes no args and can be used to clear, skip — prefer explicit clear
                            if name in ("clear_store", "clear", "reset"):
                                try:
                                    fn()
                                    print(f"[RAG] Called rag_embeddings.file_embeddings().{name}()")
                                except Exception:
                                    pass
                    # also attempt module-level cache wipe
                    for attr in ("_store", "store", "index", "chunks"):
                        if hasattr(fe_inst, attr):
                            try:
                                setattr(fe_inst, attr, None)
                            except Exception:
                                pass
                except Exception as e:
                    print(f"[RAG] rag_embeddings clearing attempt failed: {e}")
        except Exception:
            # module not present - ignore
            pass

        # --- Clear in-memory app-level caches and flags (this is what prevents the app returning context) ---
        cleared_names = []
        for nm in GLOBAL_NAMES_TO_CLEAR:
            try:
                globals()[nm] = {} if nm.lower().endswith("s") else "" if "content" in nm or "my_file" in nm else []
                cleared_names.append(nm)
            except Exception:
                try:
                    # try deletion then fallback to simple assignment
                    if nm in globals():
                        del globals()[nm]
                    globals()[nm] = {} if nm.lower().endswith("s") else "" if "content" in nm or "my_file" in nm else []
                    cleared_names.append(nm)
                except Exception:
                    pass

        # Also explicitly clear alternate names that handlers might check
        try:
            globals()['USER_FILES'] = {}
        except Exception:
            pass
        try:
            globals()['RAG_STORE'] = {}
        except Exception:
            pass
        try:
            globals()['my_file_content'] = ""
        except Exception:
            pass
        try:
            globals()['uploaded_doc_ids'] = []
        except Exception:
            pass
        try:
            globals()['rag_enabled'] = False
        except Exception:
            pass

        # Finalize: best-effort persist (will not delete files if locked)
        try:
            # try to persist on one of the clients we created earlier
            for path, had_client in tried_paths:
                if had_client:
                    try:
                        client = _client_for_path(path)
                        if client is not None and hasattr(client, "persist"):
                            try:
                                client.persist()
                            except Exception:
                                pass
                    except Exception:
                        pass
        except Exception:
            pass

        # Report back
        print(f"[RAG] Reset finished. persistent_deleted_any={persistent_deleted_any}, cleared_in_memory={cleared_names}")
        socketio.emit('state_update', {
            'type': 'reset_rag_ack' if persistent_deleted_any else 'reset_rag_partial',
            'payload': "RAG persistent store cleared." if persistent_deleted_any else "RAG in-memory context cleared; persistent store may be locked or already empty.",
            'username': user
        }, room=user)

        return


    # --- LOAD FILE ---
    if t == 'load_file':


        # Initialize user file memory
        user_files = globals().get("user_files", {})
        globals()["user_files"] = user_files
        my_file_content = ""

        # --- Validate payload ---
        if not isinstance(payload, dict):
            socketio.emit('state_update', {
                'type': 'file_error',
                'payload': "Invalid payload format for load_file",
                'username': user
            }, room=user)
            return

        filename = payload.get('filename')
        filedata_b64 = payload.get('filedata')
        mime_hint = payload.get('mime') or payload.get('type')

        if not filename or filedata_b64 is None:
            socketio.emit('state_update', {
                'type': 'file_error',
                'payload': "Missing filename or filedata",
                'username': user
            }, room=user)
            return

        # --- Decode base64 file ---
        try:
            file_bytes = base64.b64decode(filedata_b64)
        except Exception as e:
            socketio.emit('state_update', {
                'type': 'file_error',
                'payload': f'Base64 decode failed: {e}',
                'username': user
            }, room=user)
            return

        # --- Detect MIME type ---
        mime_type, _ = mimetypes.guess_type(filename)
        if not mime_type and mime_hint:
            mime_type = mime_hint
        print(f"[INFO] Detected MIME type: {mime_type}")

        # optional helpers (import at top of file)
        try:
            import chardet
        except Exception:
            chardet = None

        try:
            from werkzeug.utils import secure_filename
        except Exception:
            # fallback simple sanitizer
            def secure_filename(name):
                return re.sub(r'[^a-zA-Z0-9._-]', '_', name)

        # Ensure TMP_DIR exists
        TMP_DIR = globals().get("TMP_DIR", tempfile.gettempdir())


        def extract_text_from_any(file_bytes, filename, mime_type=None):
            """
            Robust extractor for many common file types.
            Returns: extracted text (unicode string). Prints diagnostics on error.

            Notes:
            - Treats .ino (Arduino) as plain text / source code.
            - Tries utf-8 first, then chardet (if available), then latin-1 as a last resort for text files.
            """
            text = ""
            try:
                # normalize filename + get extension
                if not filename:
                    filename = "unknown"
                # use lowercased extension(s)
                _, file_ext = os.path.splitext(filename)
                ext = file_ext.lower()  # includes leading dot, e.g. '.pdf' or '.ino'

                # --- PDF handling ---
                if (mime_type == 'application/pdf') or ext == '.pdf':
                    try:
                        import fitz  # PyMuPDF
                    except Exception as e:
                        print(f"[ERROR] PyMuPDF (fitz) not available: {e}")
                        raise

                    doc = None
                    try:
                        doc = fitz.open(stream=file_bytes, filetype="pdf")
                    except Exception:
                        try:
                            doc = fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf")
                        except Exception as e:
                            print(f"[ERROR] fitz.open failed: {e}")
                            raise

                    page_text_accum = []
                    for pnum in range(doc.page_count):
                        try:
                            page = doc.load_page(pnum)
                            ptext = page.get_text("text") or ""
                            if ptext and len(ptext.strip()) > 20:
                                page_text_accum.append(ptext)
                            else:
                                # OCR fallback
                                try:
                                    from PIL import Image
                                    import pytesseract
                                    pix = page.get_pixmap(dpi=300, alpha=False)
                                    img_bytes = pix.tobytes("png")
                                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                                    ocr_text = pytesseract.image_to_string(img)
                                    page_text_accum.append(ocr_text)
                                except Exception as e:
                                    print(f"[WARN] OCR fallback failed on PDF page {pnum}: {e}")
                        except Exception as e:
                            print(f"[WARN] error processing PDF page {pnum}: {e}")
                    text = "\n".join(page_text_accum)
                    try:
                        doc.close()
                    except Exception:
                        pass

                # --- DOCX ---
                elif ext == '.docx':
                    try:
                        from docx import Document
                        doc = Document(io.BytesIO(file_bytes))
                        paragraphs = [p.text for p in doc.paragraphs if p.text]
                        text = "\n".join(paragraphs)
                    except Exception as e:
                        print(f"[ERROR] docx extraction failed: {e}")
                        raise

                # --- PPTX ---
                elif ext == '.pptx':
                    try:
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(file_bytes))
                        lines = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text") and shape.text:
                                        lines.append(shape.text)
                                except Exception:
                                    pass
                        text = "\n".join(lines)
                    except Exception as e:
                        print(f"[ERROR] pptx extraction failed: {e}")
                        raise

                # --- XLSX/XLS ---
                elif ext in ('.xlsx', '.xls'):
                    try:
                        import openpyxl
                        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                        lines = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                rowcells = [str(c) for c in row if c is not None and str(c).strip() != ""]
                                if rowcells:
                                    lines.append(" ".join(rowcells))
                        text = "\n".join(lines)
                    except Exception as e:
                        print(f"[ERROR] excel extraction failed: {e}")
                        raise

                # --- Plain text / source code (include .ino) ---
                # Add any other code extensions you want here
                elif (
                    (mime_type and mime_type.startswith('text'))
                    or ext in ('.txt', '.md', '.py', '.js', '.html', '.htm', '.css',
                               '.json', '.c', '.cpp', '.h', '.hpp', '.java', '.sh',
                               '.bash', '.ps1', '.rb', '.go', '.rs', '.php', '.pl',
                               '.scala', '.swift', '.kt', '.kts', '.ts', '.tsx', '.jsx',
                               '.ino', '.ino.txt')  # .ino explicitly included
                ):
                    # Try utf-8 first, then chardet if available, then latin-1
                    tried_encodings = []
                    def _decode_try(enc):
                        try:
                            return file_bytes.decode(enc, errors='replace')
                        except Exception:
                            return None

                    # quick utf-8 attempt
                    tried_encodings.append('utf-8')
                    text = _decode_try('utf-8')
                    if text is None or text == "":
                        # try chardet if present
                        if chardet:
                            try:
                                detected = chardet.detect(file_bytes)
                                enc = detected.get("encoding") or "utf-8"
                                tried_encodings.append(enc)
                                text = _decode_try(enc)
                            except Exception as e:
                                print(f"[WARN] chardet detection failed: {e}")
                        # final fallback to latin-1
                        if not text:
                            tried_encodings.append('latin-1')
                            text = _decode_try('latin-1') or ""
                    # normalize line endings
                    text = text.replace('\r\n', '\n').replace('\r', '\n')

                # --- Image OCR ---
                elif (mime_type and mime_type.startswith('image')) or ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'):
                    try:
                        from PIL import Image
                        import pytesseract
                        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
                        text = pytesseract.image_to_string(img)
                    except Exception as e:
                        print(f"[ERROR] image OCR failed: {e}")
                        raise

                # --- Unsupported: write raw to tmp and return empty text (or optionally return path) ---
                else:
                    try:
                        tmp_path = os.path.join(TMP_DIR, secure_filename(filename))
                        with open(tmp_path, "wb") as wf:
                            wf.write(file_bytes)
                        print(f"[WARN] Unsupported file saved to {tmp_path}")
                        text = ""
                    except Exception as e:
                        print(f"[ERROR] Could not save unsupported file: {e}")
                        raise

            except Exception as e:
                print(f"[ERROR] extract_text_from_any failed for '{filename}': {e}")
                traceback.print_exc()

            # final cleanup & info
            if isinstance(text, bytes):
                try:
                    text = text.decode("utf-8", errors="replace")
                except Exception:
                    text = str(text)

            short = (text[:200] + "...") if len(text) > 200 else text
            print(f"[INFO] extract_text_from_any: filename={filename!r} bytes={len(file_bytes) if file_bytes is not None else 'None'} -> chars={len(text)} preview={short!r}")
            return text

        # --- Extract text ---
        file_text = extract_text_from_any(file_bytes, filename, mime_type)
        file_text = re.sub(r'\s+', ' ', file_text).strip()
        print(f"[INFO] Extracted {len(file_text)} characters from '{filename}'")

        # --- Store file per-user for queries ---
        user_files[user] = file_text

        # --- Notify frontend ---
        socketio.emit('state_update', {
            'type': 'file_info',
            'payload': {'message': f"Text from '{filename}' processed for embedding"},
            'username': user
        }, room=user)

        # --- RAG EMBEDDING ---
        try:
            import chromadb

            # ✅ FIX: use new PersistentClient if available
            try:
                if hasattr(chromadb, "PersistentClient"):
                    client = chromadb.PersistentClient(path=os.path.join(BASE_DIR, "chroma_db"))
                    print("[INFO] Using new Chroma PersistentClient")
                else:
                    from chromadb.config import Settings as ChromaSettings
                    chroma_settings = ChromaSettings(
                        chroma_db_impl="duckdb+parquet",
                        persist_directory=os.path.join(BASE_DIR, "chroma_db")
                    )
                    client = chromadb.Client(chroma_settings)
                    print("[INFO] Using legacy Chroma Client")
            except Exception as e:
                print(f"[WARN] Fallback to default Chroma client: {e}")
                client = chromadb.Client()

            # Chunk text for long files
            def chunk_text(text, max_chars=1500):
                for i in range(0, len(text), max_chars):
                    yield text[i:i+max_chars]

            safe_name = secure_filename(filename)
            doc_id = f"{user}_{safe_name}_{int(time.time())}"

            # ✅ version-safe collection get/create
            if hasattr(client, "get_or_create_collection"):
                collection = client.get_or_create_collection(name="docs")
            elif hasattr(client, "create_collection"):
                collection = client.create_collection(name="docs")
            else:
                collection = client.get_collection("docs")

            # Use robust embedding helper we added
            for i, chunk in enumerate(chunk_text(file_text)):
                vect = get_ollama_embedding(chunk)
                if vect is None:
                    print(f"[WARN] embedding failed for chunk {i}; skipping chunk")
                    continue
                embedding = vect
                try:
                    collection.add(
                        ids=[f"{doc_id}_chunk{i}"],
                        embeddings=[embedding],
                        documents=[chunk]
                    )
                except Exception:
                    # Try alternate add signature (for newer Chroma versions)
                    collection.add(
                        documents=[chunk],
                        metadatas=[{"source": filename}],
                        ids=[f"{doc_id}_chunk{i}"],
                        embeddings=[embedding]
                    )

            # ✅ persist safely
            try:
                if hasattr(client, "persist"):
                    client.persist()
                elif hasattr(collection, "persist"):
                    collection.persist()
            except Exception as e:
                print(f"[WARN] Persist skipped: {e}")

            print(f"[INFO] Embedded {len(file_text)} chars into Chroma for RAG")

            # Log event
            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": f"[file_upload] {filename}",
                "response": "[stored in RAG]",
                "models": {"embed_model": "ollama/sentence-transformers"}
            }
            log_user_query(user, entry)

            # Notify frontend success
            socketio.emit('state_update', {
                'type': 'file_info',
                'payload': {
                    'message': f"Text '{filename}' processed for RAG",
                    'doc_id': doc_id,
                    'length': len(file_text),
                    'rag_available': True,
                    'embed_source': "ollama/sentence-transformers"
                },
                'username': user
            }, room=user)

        except Exception as e:
            print(f"[ERROR] RAG embedding failed: {e}")
            socketio.emit('state_update', {
                'type': 'file_error',
                'payload': f'RAG processing failed: {e}',
                'username': user
            }, room=user)

        return


    # --- QUERY ---
    if t == 'query':
        text = payload if isinstance(payload, str) else payload.get('query', '')
        print(f"[EVENT] query from user='{user}' text='{text}'")

        if not text:
            print("[WARN] Empty query text.")
            return

        user_files = globals().get("user_files", {})
        if any(x in text.lower() for x in ["this file", "this story", "this document", "this code", "this text"]) and user in user_files:
            try:
                from rag_embeddings import file_embeddings
                fe = file_embeddings()
                fe.build_store(user_files[user])
                context_chunks = fe.retrieve(query=text, top_k=20)
                context = "\n".join(f"- {chunk.strip()}" for chunk in context_chunks) if context_chunks else "No relevant items found."

                final_prompt = f"{text}\n\nContext:\n{context}"
                print(f"[RAG QUERY] Final prompt:\n{final_prompt[:400]}...")

            except Exception as e:
                print(f"[RAG ERROR QUERY] {e}")
                final_prompt = text

            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": final_prompt,
                "response": None,
                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
            }
            log_user_query(user, entry)
            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
            return

        else:
            entry = {
                "username": user,
                "timestamp": datetime.now().isoformat(),
                "query": text,
                "response": None,
                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
            }
            log_user_query(user, entry)
            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
            return


# HTTP endpoints
@app.route('/main_responses/<username>', methods=['GET'])
def main_responses(username):
    path = os.path.join(TMP_DIR, f"{username}_main.json")
    if not os.path.exists(path):
        return jsonify([])
    with open(path, 'r', encoding='utf-8') as f:
        try:
            return jsonify(json.load(f))
        except json.JSONDecodeError:
            return jsonify([])

@app.route('/user_logs/<username>', methods=["GET"])
def get_user_logs(username):
    logs = read_logs(username)
    print(f"[DEBUG] Loaded logs for {username}: {len(logs)} entries")
    return jsonify(logs)

@app.post("/upload_audio")
def upload_audio():
    os.makedirs(TMP_DIR, exist_ok=True)
    webm_path = os.path.join(TMP_DIR, "temp_audio.webm")
    wav_path = os.path.join(TMP_DIR, "temp_audio.wav")

    if 'audio' not in request.files:
        return jsonify({"error": "No audio file"}), 400
    request.files["audio"].save(webm_path)

    # convert with ffmpeg if available
    if os.system(f'ffmpeg -y -i "{webm_path}" -ar 16000 -ac 1 -f wav "{wav_path}"') != 0:
        return jsonify({"error": "ffmpeg failed"}), 500

    try:
        text = whisper_model.transcribe(wav_path).get("text", "").strip()
        cleaned = webui_listener.listen_text(text)
    except Exception:
        cleaned = ""

    for p in (webm_path, wav_path):
        try:
            os.remove(p)
        except:
            pass

    return jsonify({"transcript": cleaned})

@app.post('/upload_file')
def upload_file_http():
    # HTTP multipart upload alternative to socket-based upload
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    username = request.form.get('username', 'Tjaart')
    filename = secure_filename(file.filename)
    tmp_path = os.path.join(TMP_DIR, filename)
    file.save(tmp_path)

    # Try to detect type and if text, process into Chroma
    mime_type, _ = mimetypes.guess_type(filename)

    if mime_type and mime_type.startswith('image'):
        print(f"[HTTP FILE] Image uploaded: {filename} by {username}")
        socketio.emit('state_update', {'type': 'file_info', 'payload': f"Image '{filename}' received", 'username': username}, room=username)
        return jsonify({"status": "ok", "type": "image"})

    if mime_type and (mime_type.startswith('text') or filename.lower().endswith(('.txt', '.md', '.py'))):
        print(f"[HTTP FILE] Text uploaded: {filename} by {username}")
        try:
            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                file_text = f.read()
        except Exception as e:
            return jsonify({"error": f"Could not read file: {e}"}), 500

        # Same RAG logic as socket handler
        try:
            client = None
            if chromadb is not None:
                try:
                    from chromadb.config import Settings as ChromaSettings
                    chroma_settings = ChromaSettings(
                        chroma_db_impl="duckdb+parquet",
                        persist_directory=os.path.join(BASE_DIR, "chroma_db")
                    )
                    client = chromadb.Client(chroma_settings)
                except Exception:
                    try:
                        client = chromadb.Client()
                    except Exception:
                        client = None

            embedding = None
            emb_source = None
            try:
                # Use robust helper to create embedding
                embedding = get_ollama_embedding(file_text)
                if embedding is not None:
                    emb_source = 'ollama'
            except Exception:
                embedding = None
                emb_source = None

            if embedding is None:
                try:
                    # fallback - sentence-transformers handled in helper already but try again gracefully
                    from sentence_transformers import SentenceTransformer
                    st_model = SentenceTransformer('all-MiniLM-L6-v2')
                    vect = st_model.encode(file_text)
                    embedding = vect.tolist() if hasattr(vect, 'tolist') else list(vect)
                    emb_source = 'sentence-transformers'
                except Exception:
                    embedding = None
                    emb_source = None

            doc_id = f"{username}_{filename}_{int(time.time())}"
            if embedding is not None and client is not None:
                try:
                    collection = client.get_or_create_collection(name="docs")
                    collection.add(ids=[doc_id], embeddings=[embedding], documents=[file_text])
                    try:
                        client.persist()
                    except Exception:
                        pass
                except Exception:
                    client = None

            if embedding is None or client is None:
                fallback_index = os.path.join(TMP_DIR, "local_rag_index.json")
                try:
                    if os.path.exists(fallback_index):
                        with open(fallback_index, 'r', encoding='utf-8') as fh:
                            idx = json.load(fh)
                    else:
                        idx = []
                except Exception:
                    idx = []
                idx_entry = {
                    'doc_id': doc_id,
                    'username': username,
                    'filename': filename,
                    'path': tmp_path,
                    'length': len(file_text),
                    'embedding_source': emb_source,
                    'timestamp': datetime.now().isoformat()
                }
                idx.append(idx_entry)
                try:
                    with open(fallback_index, 'w', encoding='utf-8') as fh:
                        json.dump(idx, fh, ensure_ascii=False, indent=2)
                except Exception:
                    pass

            log_user_query(username, {"username": username, "timestamp": datetime.now().isoformat(), "query": f"[file_upload] {filename}", "response": "[stored in RAG]" if embedding is not None else "[saved - RAG unavailable]", "models": {"embed_model": emb_source or "none"}})

            # Provide frontend-friendly summary (avoid undefined vars)
            summary_prompt = f"[DOCUMENT STORED] filename={filename}, doc_id={doc_id}, length={len(file_text)}, embed_source={emb_source or 'none'}"
            socketio.emit('state_update', {
                'type': 'query',
                'payload': summary_prompt,
                'username': username
            }, room=username)

            return jsonify({"status": "ok", "type": "text", "doc_id": doc_id})
        except Exception as e:
            print(f"[ERROR] HTTP RAG failed: {e}")
            return jsonify({"error": str(e)}), 500

    return jsonify({"status": "ok", "type": "unknown"})


@app.route('/edge_tts', methods=['POST'])
def edge_tts_endpoint():
    """
    Synthesize text with edge-tts and return MP3 bytes.
    Frontend posts JSON: { text, voice?, style? }.
    Returns: audio/mpeg binary (200) or JSON error (non-200).
    """
    data = request.get_json(silent=True) or {}
    text = data.get('text') or data.get('response') or ""
    voice = data.get('voice', 'en-US-GuyNeural')
    style = data.get('style', None)  # currently unused in this simple example

    if not text:
        return jsonify({"error": "No 'text' provided"}), 400

    # Try importing edge-tts
    try:
        import edge_tts
    except Exception as e:
        return jsonify({"error": f"edge-tts not installed or import failed: {e}"}), 500

    # create a temporary file for the mp3
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
    tmp_path = tmp.name
    tmp.close()

    async def synth_to_file(txt, voice_name):
        # Use simple text; if you want SSML or style-based SSML, build it here.
        comm = edge_tts.Communicate(txt, voice_name)
        # save will write mp3 to path
        await comm.save(tmp_path)

    try:
        # Run the async synth (blocks until complete)
        asyncio.run(synth_to_file(text, voice))
    except Exception as e:
        # cleanup
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
        return jsonify({"error": f"edge-tts synthesis failed: {e}"}), 500

    # Return the mp3 file bytes
    try:
        # send_file will set appropriate headers
        resp = send_file(tmp_path, mimetype='audio/mpeg', as_attachment=False)
        # Allow cross-origin if needed (your SocketIO allowed * already)
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp
    finally:
        # cleanup file after response has been scheduled (Flask may still be reading it,
        # but removing the file here is okay on many OSes; if you prefer, delete in a worker)
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

@app.route('/login', methods=["POST"])
def login_user():
    data = request.json
    success, msg = do_login(data["username"], data["password"])
    return jsonify({"success": success, "message": msg})

@app.route('/signup', methods=["POST"])
def signup_user():
    data = request.json
    success, msg = do_signup(data["username"], data["password"])
    return jsonify({"success": success, "message": msg})

@app.route('/users', methods=["GET"]) 
def list_users():
    return jsonify(get_users())

@app.route('/delete_user', methods=["POST"]) 
def remove_user():
    data = request.json
    success = delete_user(data["username"])
    return jsonify({"success": success})

@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def serve(path):
    if path != "" and os.path.exists(os.path.join(app.static_folder, path)):
        return send_from_directory(app.static_folder, path)
    return send_from_directory(app.static_folder, 'index.html')

if __name__ == '__main__':
    context = ('cert.pem', 'key.pem') if os.path.exists('cert.pem') and os.path.exists('key.pem') else None
    socketio.run(app, host='0.0.0.0', port=5001, ssl_context=context)












######import ollama
######import re
######import os
######import json
######import time
######import base64
######import mimetypes
######import tempfile
######from datetime import datetime
######from werkzeug.utils import secure_filename
######
######import whisper
######from flask import Flask, request, jsonify, send_from_directory
######from flask_socketio import SocketIO, emit, join_room, leave_room
######
####### Optional libs (may not be installed in every env)
######try:
######    import chromadb
######except Exception:
######    chromadb = None
######try:
######    import ollama
######except Exception:
######    ollama = None
######
######from listenWEBUI import WEBUIListenModule
######from auth_manager import login as do_login, signup as do_signup, get_users, delete_user
######from query_logger import log_user_query, log_user_response, read_logs
######
######import chardet
######from pdfminer.high_level import extract_text
######from pdf2image import convert_from_path
######import pytesseract
######import base64
######import io
######import PyPDF2
######from PyPDF2 import PdfFileReader
######import pdfplumber
######from langchain_chroma import Chroma
######
######import asyncio
######import tempfile
######from flask import send_file   # modify existing import line to include send_file if needed
######
######import io, re, base64, mimetypes, tempfile, os, time
######import chardet
######from werkzeug.utils import secure_filename
######
######import io, re, base64, mimetypes, tempfile, os, time
######import chardet
######from werkzeug.utils import secure_filename
######from datetime import datetime
######        
######BASE_DIR = os.path.dirname(os.path.abspath(__file__))
######TMP_DIR = os.path.join(BASE_DIR, 'tmp')
######os.makedirs(TMP_DIR, exist_ok=True)
######
####### Initialize listener and Whisper
######webui_listener = WEBUIListenModule()
######whisper_model = whisper.load_model("base.en")
######
######app = Flask(__name__, static_folder='./webui-src/dist', static_url_path='')
######socketio = SocketIO(app, cors_allowed_origins='*', async_mode='threading')
######
######
######chromadb.api.client.SharedSystemClient.telemetry_enabled = False
######
######
####### Module-level state
######state = {
######    'logs': [],
######    'queries': [],
######    'resp': [],
######    'settings': {
######        'use_whisper': False,
######        'use_vosk': True,
######        'enter_submits': True
######    }
######}
######
####### Session tracking
######session_users = {}  # Maps SID -> username
######
######my_file_content = ""
######
######collection = None   # Chroma collection shared across events
######os.environ["CHROMA_TELEMETRY_ENABLED"] = "false"
######
########from langchain.text_splitter import RecursiveCharacterTextSplitter
######import PyPDF2
######
######
######model_selected = ""
######thinkbot_model = ""
######rag_clear = False
######detected_model = ""
######RAG_model = ""
######
######import io
######import os
######import tempfile
######import traceback
######import base64
######import re
######
######
######
######
####### ======= Query helper (can be called from elsewhere in backend) =======
######def query_with_retrieval(prompt: str, n_results: int = 1):
######    """Retrieve top chunks from Chroma and run generation with Ollama."""
######    try:
######        if ollama is None or chromadb is None:
######            raise RuntimeError("Ollama or Chroma not available")
######
######        q_emb = ollama.embeddings(model="mxbai-embed-large", prompt=prompt)
######        q_vector = q_emb.get("embedding") if isinstance(q_emb, dict) else q_emb
######        client = chromadb.Client()
######        collection = client.get_or_create_collection(name="docs")
######        res = collection.query(query_embeddings=[q_vector], n_results=n_results)
######        top_doc = res.get("documents", [[]])[0][0] if res.get("documents") else ""
######    except Exception as e:
######        print(f"[WARN] retrieval failed: {e}")
######        top_doc = ""
######
######    gen_prompt = f"Using this data: {top_doc}\n\nRespond to: {prompt}"
######    try:
######        out = ollama.generate(model="tinyllama", prompt=gen_prompt, stream=False)
######        text = out.get("response") or out.get("text") or str(out)
######    except Exception as e:
######        text = f"[generation failed: {e}]"
######    return text
######
######def extract_pdf_text_with_ocr(pdf_path):
######    try:
######        text = extract_text(pdf_path)
######        if len(text.strip()) > 10:  # if decent text extracted, return it
######            return text
######        else:
######            # Fallback to OCR if text is too short or empty
######            print("[INFO] PDF text extraction empty or too short, using OCR fallback.")
######            pages = convert_from_path(pdf_path)
######            ocr_text = ""
######            for page in pages:
######                ocr_text += pytesseract.image_to_string(page)
######            return ocr_text
######    except Exception as e:
######        print(f"[ERROR] PDF extraction + OCR failed: {e}")
######        return ""
######
######
####### Utility: safe save
######def _save_tmp(filename, data_bytes):
######    safe = secure_filename(filename)
######    path = os.path.join(TMP_DIR, safe)
######    with open(path, 'wb') as f:
######        f.write(data_bytes)
######    return path
######
####### Socket handlers
######@socketio.on('connect')
######def connect():
######    emit('full_state', state)
######
######@socketio.on('disconnect')
######def disconnect():
######    sid = request.sid
######    user = session_users.pop(sid, None)
######    if user:
######        leave_room(user)
######        print(f"[DISCONNECT] {sid} left room {user}")
######
######
####### --- helper: singleton chroma client creator (safe) ---
######
######_CHROMA_CLIENT = None
######
######def get_chroma_client(persist_directory=None, chroma_db_impl="duckdb+parquet"):
######    global _CHROMA_CLIENT
######    if _CHROMA_CLIENT is not None:
######        return _CHROMA_CLIENT
######    try:
######        import chromadb
######        from chromadb.config import Settings as ChromaSettings
######    except Exception as e:
######        print(f"[RAG] chromadb import failed: {e}")
######        return None
######
######    try:
######        settings = ChromaSettings(chroma_db_impl=chroma_db_impl,
######                                  persist_directory=persist_directory)
######
######        client = None
######        try:
########            from chromadb.config import Settings as ChromaSettings
######            chroma_settings = ChromaSettings(
######                chroma_db_impl="duckdb+parquet",
######                persist_directory=os.path.join(BASE_DIR, "chroma_db")
######            )
######            client = chromadb.Client(chroma_settings)
######        except Exception:
######            # fallback to default client constructor; may still fail, handle below
######            try:
######                client = chromadb.Client()
######            except Exception as e:
######                client = None
######                print(f"[RAG] chromadb.Client fallback failed: {e}")
######
######
######            if client is None:
######                client = get_chroma_client(persist_directory=CHROMA_DIR)
######
######        return client
######
######
######
######
########        _CHROMA_CLIENT = chromadb.Client(settings=settings)
########        print(f"[RAG] Chroma client created with persist_directory={persist_directory}")
########        return _CHROMA_CLIENT
######    except Exception as e:
######        print(f"[RAG] Embedded Chroma init failed: {e}; trying fallback chromadb.Client()")
######        try:
######            _CHROMA_CLIENT = chromadb.Client()
######            print("[RAG] Chroma client created via fallback chromadb.Client()")
######            return _CHROMA_CLIENT
######        except Exception as e2:
######            print(f"[RAG] chromadb.Client() fallback failed: {e2}")
######            _CHROMA_CLIENT = None
######            return None
######
######_MODEL_STR_PATTERN = re.compile(r"^[A-Za-z0-9_\-]+(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)?$")
######
######def extract_model(payload) -> str:
######    """
######    Robustly extract and clean a model name. Always returns a string ("" if none found).
######
######    Rules:
######      - If payload is a dict and contains thinkbot_model/model/model_used -> use it.
######      - Else if payload is a dict and has 'payload' whose value is a dict -> descend and try again.
######      - Else if payload is a plain string and *looks like a model name* (pattern) -> accept it.
######      - Otherwise return "".
######
######    Example results:
######      {'payload': {'thinkbot_model': 'moondream'}} -> "moondream"
######      {'thinkbot_model': 'deepseek-r1-1.5b'}    -> "deepseek-r1"
######      "deepseek-r1"                             -> "deepseek-r1"
######      {'type': 'log', 'payload': 'I see, a man ...'} -> ""
######    """
######    model_raw = ""
######
######    # If dict, prefer direct model keys at top-level
######    if isinstance(payload, dict):
######        # Direct model keys first
######        if "detected_model" in payload or "model" in payload or "model_used" in payload:
######            model_raw = payload.get("detected_model") or payload.get("model") or payload.get("model_used") or ""
######        else:
######            # Only descend into 'payload' if it's a dict (avoid descending into plain log strings)
######            inner = payload.get("payload")
######            if isinstance(inner, dict):
######                model_raw = inner.get("detected_model") or inner.get("model") or inner.get("model_used") or ""
######            else:
######                # No model found in this event
######                return ""
######    else:
######        # payload not a dict; if it's a plain string that *looks like a model* accept it
######        if isinstance(payload, str) and _MODEL_STR_PATTERN.match(payload.strip()):
######            model_raw = payload.strip()
######        else:
######            return ""
######
######    model_raw = (model_raw or "").strip()
######    if not model_raw:
######        return ""
######
######    # Clean suffixes like ":1.5b", "-1.5b", ":8b", "-8b", etc.
######    model_clean = re.sub(r"(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)$", "", model_raw)
######
######    return model_clean
######
######
######@socketio.on('gui_event')
######def handle_event(data):
######    global client, collection, detected_model, thinkbot_model, model_selected, rag_clear  # Declare globals at the very top
######    sid = request.sid
######    print(f"[DEBUG] Received gui_event from sid={sid}: {data!r}")
######
######    user = data.get('username') if isinstance(data, dict) else None
######    t = data.get('type') if isinstance(data, dict) else None
######    payload = data.get('payload') if isinstance(data, dict) else None
######
######    event = data
######    print("[EVENT] Received event 'new event' : ", str(event))
######
######    model_used = extract_model(event)
######    print(f"[EVENT] Received gui_event 'model_used' : {model_used}")
######
######    if not user:
######        user = 'Tjaart'  # fallback
######    # Track user per session
######    session_users[sid] = user
######    join_room(user)
######    global my_file_content
######    
######    # --- LOGIN ---
######    if t == 'login':
######        print(f"[LOGIN] User: {user}")
######        socketio.emit('state_update', {
######            'type': 'login_ack',
######            'username': user
######        }, room=user)
######        return
######
######    # --- LOG ---
######    if t == 'log':
######        # Always produce a structured entry so frontend can read models reliably
######        if isinstance(payload, str):
######            # Keep old logging side-effect
######            log_user_response(user, payload)
######
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": "",
######                "response": payload,
######                "models": {"thinkbot": model_used or ""}
######            }
######
######            socketio.emit('state_update', {
######                'type': 'log',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######        elif isinstance(payload, dict):
######            # Ensure models is a dict and add thinkbot
######            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
######            models = dict(models)  # copy-safe
######            models["thinkbot"] = model_used or ""
######            query = payload.get("query", "")
######
######            if query:
######                guery = query.replace(":","")
######                print(f"New Query supplied : {query}")
######            else:
######                print("There is no query supplied...")
######
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": guery,
######                "response": payload.get("response", ""),
######                "models": models
######            }
######
######            log_user_query(user, entry)
######
######            socketio.emit('state_update', {
######                'type': 'log',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######            socketio.emit('state_update', {
######                'type': 'query',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######        return
######
######    # --- RESP ---
######    if t == 'resp':
######        if isinstance(payload, str):
######            log_user_response(user, payload)
######
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": "",
######                "response": payload,
######                "models": {"thinkbot": model_used or ""}
######            }
######
######            socketio.emit('state_update', {
######                'type': 'log',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######        elif isinstance(payload, dict):
######            # Preserve any existing models and inject thinkbot
######            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
######            models = dict(models)
######            models["thinkbot"] = model_used or ""
######            query = payload.get("query", "")
######            
######            if query:
######                query = query.replace(":","")
######                print(f"New Query supplied : {query}")
######            else:
######                print("No Query supplied...")
######                 
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": query,
######                "response": payload.get("response", ""),
######                "models": models
######            }
######
######            log_user_query(user, entry)
######
######            socketio.emit('state_update', {
######                'type': 'log',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######            socketio.emit('state_update', {
######                'type': 'query',
######                'payload': entry,
######                'username': user
######            }, room=user)
######
######            print(f" \n thinkbot_model : {model_used} \n")
######
######        return
######
######    # --- SETTING ---
######    if t == 'setting':
######        socketio.emit('state_update', {
######            'type': 'setting',
######            'payload': payload,
######            'username': user
######        }, room=user)
######        return
######
######
######    if t in ('setting', 'detected_model'):
######        data_payload = data.get('payload')
######        user = data.get('username') or 'default_user'
########        new_detected_model = data.get('detected_model')
######
######        # Try to extract the thinkbot model name from different payload shapes:
######        model_selected = None
######        
######        if isinstance(data_payload, dict):
######            # common: payload = { 'detected_model': 'tinyllama' }
######            if 'detected_model' in data_payload:
######                model_selected = data_payload['detected_model']
######            else:
######                # fallback: payload may be { '<some_model_key>': '<name>' }
######                # find any key that ends with '_model' and pick its value
######                model_keys = [k for k in data_payload.keys() if k.endswith('_model')]
######                if len(model_keys) == 1:
######                    model_selected = data_payload[model_keys[0]]
######        elif isinstance(data_payload, str):
######            # sometimes frontend sends a raw string payload (e.g. from 'think_model' emit)
######            model_selected = data_payload
######
######        # If we found a model, assign and notify
######        if model_selected:
######            # ensure we set the module/global variable (declare `global detected_model` above if needed)
######            try:
######                global detected_model
######            except NameError:
######                # if detected_model doesn't exist yet, create it
######                detected_model = None
######
######            detected_model = model_selected  # correct assignment (not ==)
######            print(f"\n The detected_model for RAG is '{detected_model} \n'")
######            print(f"\n The new detected_model SETTINGS is now '{detected_model}' \n")
######
######            # Emit a consistent state_update so clients can sync
######            socketio.emit('state_update', {
######                'type': 'setting',
######                'payload': {'detected_model': detected_model},
######                'username': user
######            }, room=user)
######        else:
######            # No model found — still forward the raw payload for other settings
######            print(f"[setting] received payload but no model detected: {data_payload}")
######            socketio.emit('state_update', {
######                'type': 'setting',
######                'payload': data_payload,
######                'username': user
######            }, room=user)
######
######        return
######
######
######    # --- SETTING ---
######    if t == 'think_model':
######
######        payload = data.get('payload')
######        user = data.get('username') or 'default_user'
######        model_selected = data.get('think_model')
######
######        print(f"The model for RAG 'model_selected' is {model_selected}")
######        return
######
######
######    # --- RESET RAG ---
######    if t == 'reset_rag':
######        import getpass, traceback, time
######
######        # globals we will clear (try many names used across your app)
######        GLOBAL_NAMES_TO_CLEAR = [
######            'RAG_STORE', 'rag_store', 'USER_FILES', 'user_files',
######            'my_file_content', 'my_file_contents', 'uploaded_doc_ids',
######            'rag_enabled', '_CHROMA_CLIENT', '_CHROMA_CLIENT_SETTINGS', 'collection'
######        ]
######
######        CHROMA_DIRS_TO_TRY = []
######        # primary location (your BASE_DIR)
######        CHROMA_DIRS_TO_TRY.append(os.path.join(BASE_DIR, "chroma_db"))
######        # alternate locations you used in other codepaths
######        CHROMA_DIRS_TO_TRY.append(os.path.join(os.getcwd(), "chroma_db"))
######        CHROMA_DIRS_TO_TRY.append(os.path.join(os.getcwd(), "chroma_store"))
######        # module-local chroma_db if running from modules folder
######        CHROMA_DIRS_TO_TRY.append(os.path.join(os.path.dirname(__file__), "chroma_db"))
######        # dedupe
######        CHROMA_DIRS_TO_TRY = [os.path.normpath(p) for p in dict.fromkeys([p for p in CHROMA_DIRS_TO_TRY if p])]
######
######        user = data.get('username') if isinstance(data, dict) else user
######        user = user or getpass.getuser()
######        print(f"[RAG] reset_rag requested by {user}")
######        print(f"[RAG] CHROMA_DIRS to try: {CHROMA_DIRS_TO_TRY}")
######
######        # Attempt to import chromadb and silence telemetry to reduce noisy logs
######        chromadb_mod = globals().get('chromadb', None)
######        try:
######            if chromadb_mod is None:
######                import chromadb as chromadb_imported
######                chromadb_mod = chromadb_imported
######                globals()['chromadb'] = chromadb_mod
######            try:
######                setattr(chromadb_mod, "telemetry_enabled", False)
######            except Exception:
######                pass
######            try:
######                telemetry = getattr(chromadb_mod, "telemetry", None)
######                if telemetry is not None:
######                    for name in ("capture", "capture_event", "send_event", "capture_exception"):
######                        if hasattr(telemetry, name):
######                            try:
######                                setattr(telemetry, name, lambda *a, **kw: None)
######                            except Exception:
######                                pass
######            except Exception:
######                pass
######        except Exception as e:
######            print(f"[RAG] chromadb import/telemetry patch failed: {e}")
######            chromadb_mod = None
######
######        # helper: attempt version-tolerant deletion on a client/collection
######        def _try_delete_collection_version_tolerant(client, colname="docs"):
######            """
######            Returns True if we believe persistent deletion succeeded (or store is effectively empty).
######            """
######            if client is None:
######                return False
######            deleted = False
######            try:
######                # preferred high-level delete_collection
######                if hasattr(client, "delete_collection"):
######                    try:
######                        # some versions accept positional, some keyword
######                        try:
######                            client.delete_collection(colname)
######                        except TypeError:
######                            client.delete_collection(name=colname)
######                        print(f"[RAG] client.delete_collection('{colname}') succeeded")
######                        return True
######                    except Exception as exc:
######                        print(f"[RAG] client.delete_collection failed: {exc!r}")
######
######                # obtain collection object in a version tolerant way
######                coll = None
######                try:
######                    if hasattr(client, "get_or_create_collection"):
######                        coll = client.get_or_create_collection(name=colname)
######                    else:
######                        # try get_collection then create_collection
######                        try:
######                            if hasattr(client, "get_collection"):
######                                coll = client.get_collection(colname)
######                        except Exception:
######                            coll = None
######                        if coll is None and hasattr(client, "create_collection"):
######                            coll = client.create_collection(colname)
######                except Exception as e:
######                    print(f"[RAG] failed to obtain collection object: {e}")
######                    coll = None
######
######                if coll is None:
######                    # nothing more we can do on this client
######                    return False
######
######                # 1) try no-arg delete (supported on some versions)
######                try:
######                    if hasattr(coll, "delete"):
######                        try:
######                            coll.delete()
######                            print("[RAG] collection.delete() (no-arg) invoked")
######                            return True
######                        except ValueError as ve:
######                            # API demands ids/where; fallback below
######                            print(f"[RAG] collection.delete() raised ValueError (expected on some APIs): {ve}")
######                        except Exception as exc:
######                            print(f"[RAG] collection.delete() (no-arg) failed: {exc!r}")
######                except Exception:
######                    pass
######
######                # 2) list ids then delete by ids
######                try:
######                    ids = []
######                    try:
######                        res = coll.get(include=['ids'])
######                        if isinstance(res, dict):
######                            ids_field = res.get('ids', [[]])
######                        else:
######                            ids_field = []
######                    except Exception:
######                        try:
######                            res = coll.get()
######                            ids_field = res.get('ids', [[]]) if isinstance(res, dict) else []
######                        except Exception:
######                            ids_field = []
######
######                    # Normalize shapes like [[id1,id2]] or [id1,id2]
######                    if isinstance(ids_field, list) and len(ids_field) > 0 and isinstance(ids_field[0], list):
######                        ids = ids_field[0]
######                    elif isinstance(ids_field, list):
######                        ids = ids_field
######                    ids = [str(i) for i in ids if i is not None]
######
######                    if ids:
######                        try:
######                            coll.delete(ids=ids)
######                            print(f"[RAG] collection.delete(ids=...) succeeded ({len(ids)} ids)")
######                            return True
######                        except Exception as exc:
######                            print(f"[RAG] collection.delete(ids=...) failed: {exc!r}")
######                    else:
######                        # no ids -> treat as empty
######                        print("[RAG] collection.get() returned no ids; treating store as empty")
######                        return True
######                except Exception as exc:
######                    print(f"[RAG] listing+delete attempt failed: {exc!r}")
######
######                # 3) try deleting by metadata 'where' shapes (best-effort)
######                try:
######                    # try a few shapes that different chroma versions expect
######                    for where_shape in (
######                        {"username": user},
######                        {"metadata": {"username": user}},
######                        {"metadatas": {"username": user}},
######                        {"$or": [{"metadata.username": user}, {"username": user}]},
######                    ):
######                        try:
######                            coll.delete(where=where_shape)
######                            print(f"[RAG] collection.delete(where={where_shape}) invoked")
######                            return True
######                        except Exception:
######                            continue
######                except Exception:
######                    pass
######
######            except Exception as e:
######                print(f"[RAG] _try_delete_collection_version_tolerant error: {e}\n{traceback.format_exc()}")
######
######            return deleted
######
######        # helper to create client for a given persist_directory (try to reuse get_chroma_client helper if available)
######        def _client_for_path(persist_path):
######            # prefer your helper if present
######            try:
######                helper = globals().get("get_chroma_client")
######                if callable(helper):
######                    res = helper(persist_directory=persist_path)
######                    # helper may return (client, path)
######                    if isinstance(res, tuple) and len(res) >= 1:
######                        return res[0]
######                    return res
######            except Exception as e:
######                print(f"[RAG] get_chroma_client helper failed for {persist_path}: {e}")
######
######            # fallback direct creation
######            try:
######                import chromadb
######                try:
######                    from chromadb.config import Settings as ChromaSettings
######                    settings = ChromaSettings(chroma_db_impl="duckdb+parquet", persist_directory=persist_path)
######                    client = chromadb.Client(settings)
######                    return client
######                except Exception:
######                    try:
######                        client = chromadb.Client()
######                        return client
######                    except Exception as e2:
######                        print(f"[RAG] chromadb.Client() fallback failed for {persist_path}: {e2}")
######                        return None
######            except Exception as e:
######                print(f"[RAG] chromadb import failed in _client_for_path: {e}")
######                return None
######
######        persistent_deleted_any = False
######        tried_paths = []
######        for path in CHROMA_DIRS_TO_TRY:
######            try:
######                client = _client_for_path(path)
######                tried_paths.append((path, bool(client)))
######                if client is None:
######                    print(f"[RAG] No client for {path}; skipping persistent deletion there.")
######                    continue
######                ok = _try_delete_collection_version_tolerant(client, colname="docs")
######                if ok:
######                    print(f"[RAG] Persistent deletion succeeded for path {path}")
######                    persistent_deleted_any = True
######                else:
######                    print(f"[RAG] Persistent deletion did not succeed for path {path} (no error, may be empty or locked).")
######            except Exception as e:
######                print(f"[RAG] exception trying path {path}: {e}\n{traceback.format_exc()}")
######
######        # Try to clear rag_embeddings module caches if present (best-effort)
######        try:
######            import rag_embeddings
######            FECls = getattr(rag_embeddings, "file_embeddings", None)
######            if FECls:
######                try:
######                    fe_inst = FECls()
######                    # try well-named clear methods
######                    for name in ("clear_store", "clear", "reset", "build_store"):
######                        fn = getattr(fe_inst, name, None)
######                        if callable(fn):
######                            # if build_store takes no args and can be used to clear, skip — prefer explicit clear
######                            if name in ("clear_store", "clear", "reset"):
######                                try:
######                                    fn()
######                                    print(f"[RAG] Called rag_embeddings.file_embeddings().{name}()")
######                                except Exception:
######                                    pass
######                    # also attempt module-level cache wipe
######                    for attr in ("_store", "store", "index", "chunks"):
######                        if hasattr(fe_inst, attr):
######                            try:
######                                setattr(fe_inst, attr, None)
######                            except Exception:
######                                pass
######                except Exception as e:
######                    print(f"[RAG] rag_embeddings clearing attempt failed: {e}")
######        except Exception:
######            # module not present - ignore
######            pass
######
######        # --- Clear in-memory app-level caches and flags (this is what prevents the app returning context) ---
######        cleared_names = []
######        for nm in GLOBAL_NAMES_TO_CLEAR:
######            try:
######                globals()[nm] = {} if nm.lower().endswith("s") else "" if "content" in nm or "my_file" in nm else []
######                cleared_names.append(nm)
######            except Exception:
######                try:
######                    # try deletion then fallback to simple assignment
######                    if nm in globals():
######                        del globals()[nm]
######                    globals()[nm] = {} if nm.lower().endswith("s") else "" if "content" in nm or "my_file" in nm else []
######                    cleared_names.append(nm)
######                except Exception:
######                    pass
######
######        # Also explicitly clear alternate names that handlers might check
######        try:
######            globals()['USER_FILES'] = {}
######        except Exception:
######            pass
######        try:
######            globals()['RAG_STORE'] = {}
######        except Exception:
######            pass
######        try:
######            globals()['my_file_content'] = ""
######        except Exception:
######            pass
######        try:
######            globals()['uploaded_doc_ids'] = []
######        except Exception:
######            pass
######        try:
######            globals()['rag_enabled'] = False
######        except Exception:
######            pass
######
######        # Finalize: best-effort persist (will not delete files if locked)
######        try:
######            # try to persist on one of the clients we created earlier
######            for path, had_client in tried_paths:
######                if had_client:
######                    try:
######                        client = _client_for_path(path)
######                        if client is not None and hasattr(client, "persist"):
######                            try:
######                                client.persist()
######                            except Exception:
######                                pass
######                    except Exception:
######                        pass
######        except Exception:
######            pass
######
######        # Report back
######        print(f"[RAG] Reset finished. persistent_deleted_any={persistent_deleted_any}, cleared_in_memory={cleared_names}")
######        socketio.emit('state_update', {
######            'type': 'reset_rag_ack' if persistent_deleted_any else 'reset_rag_partial',
######            'payload': "RAG persistent store cleared." if persistent_deleted_any else "RAG in-memory context cleared; persistent store may be locked or already empty.",
######            'username': user
######        }, room=user)
######
######        return
######
######
######    # --- LOAD FILE ---
######    if t == 'load_file':
######
######
######        # Initialize user file memory
######        user_files = globals().get("user_files", {})
######        globals()["user_files"] = user_files
######        my_file_content = ""
######
######        # --- Validate payload ---
######        if not isinstance(payload, dict):
######            socketio.emit('state_update', {
######                'type': 'file_error',
######                'payload': "Invalid payload format for load_file",
######                'username': user
######            }, room=user)
######            return
######
######        filename = payload.get('filename')
######        filedata_b64 = payload.get('filedata')
######        mime_hint = payload.get('mime') or payload.get('type')
######
######        if not filename or filedata_b64 is None:
######            socketio.emit('state_update', {
######                'type': 'file_error',
######                'payload': "Missing filename or filedata",
######                'username': user
######            }, room=user)
######            return
######
######        # --- Decode base64 file ---
######        try:
######            file_bytes = base64.b64decode(filedata_b64)
######        except Exception as e:
######            socketio.emit('state_update', {
######                'type': 'file_error',
######                'payload': f'Base64 decode failed: {e}',
######                'username': user
######            }, room=user)
######            return
######
######        # --- Detect MIME type ---
######        mime_type, _ = mimetypes.guess_type(filename)
######        if not mime_type and mime_hint:
######            mime_type = mime_hint
######        print(f"[INFO] Detected MIME type: {mime_type}")
######
######        # optional helpers (import at top of file)
######        try:
######            import chardet
######        except Exception:
######            chardet = None
######
######        try:
######            from werkzeug.utils import secure_filename
######        except Exception:
######            # fallback simple sanitizer
######            def secure_filename(name):
######                return re.sub(r'[^a-zA-Z0-9._-]', '_', name)
######
######        # Ensure TMP_DIR exists
######        TMP_DIR = globals().get("TMP_DIR", tempfile.gettempdir())
######
######
######        def extract_text_from_any(file_bytes, filename, mime_type=None):
######            """
######            Robust extractor for many common file types.
######            Returns: extracted text (unicode string). Prints diagnostics on error.
######
######            Notes:
######            - Treats .ino (Arduino) as plain text / source code.
######            - Tries utf-8 first, then chardet (if available), then latin-1 as a last resort for text files.
######            """
######            text = ""
######            try:
######                # normalize filename + get extension
######                if not filename:
######                    filename = "unknown"
######                # use lowercased extension(s)
######                _, file_ext = os.path.splitext(filename)
######                ext = file_ext.lower()  # includes leading dot, e.g. '.pdf' or '.ino'
######
######                # --- PDF handling ---
######                if (mime_type == 'application/pdf') or ext == '.pdf':
######                    try:
######                        import fitz  # PyMuPDF
######                    except Exception as e:
######                        print(f"[ERROR] PyMuPDF (fitz) not available: {e}")
######                        raise
######
######                    doc = None
######                    try:
######                        doc = fitz.open(stream=file_bytes, filetype="pdf")
######                    except Exception:
######                        try:
######                            doc = fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf")
######                        except Exception as e:
######                            print(f"[ERROR] fitz.open failed: {e}")
######                            raise
######
######                    page_text_accum = []
######                    for pnum in range(doc.page_count):
######                        try:
######                            page = doc.load_page(pnum)
######                            ptext = page.get_text("text") or ""
######                            if ptext and len(ptext.strip()) > 20:
######                                page_text_accum.append(ptext)
######                            else:
######                                # OCR fallback
######                                try:
######                                    from PIL import Image
######                                    import pytesseract
######                                    pix = page.get_pixmap(dpi=300, alpha=False)
######                                    img_bytes = pix.tobytes("png")
######                                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
######                                    ocr_text = pytesseract.image_to_string(img)
######                                    page_text_accum.append(ocr_text)
######                                except Exception as e:
######                                    print(f"[WARN] OCR fallback failed on PDF page {pnum}: {e}")
######                        except Exception as e:
######                            print(f"[WARN] error processing PDF page {pnum}: {e}")
######                    text = "\n".join(page_text_accum)
######                    try:
######                        doc.close()
######                    except Exception:
######                        pass
######
######                # --- DOCX ---
######                elif ext == '.docx':
######                    try:
######                        from docx import Document
######                        doc = Document(io.BytesIO(file_bytes))
######                        paragraphs = [p.text for p in doc.paragraphs if p.text]
######                        text = "\n".join(paragraphs)
######                    except Exception as e:
######                        print(f"[ERROR] docx extraction failed: {e}")
######                        raise
######
######                # --- PPTX ---
######                elif ext == '.pptx':
######                    try:
######                        from pptx import Presentation
######                        prs = Presentation(io.BytesIO(file_bytes))
######                        lines = []
######                        for slide in prs.slides:
######                            for shape in slide.shapes:
######                                try:
######                                    if hasattr(shape, "text") and shape.text:
######                                        lines.append(shape.text)
######                                except Exception:
######                                    pass
######                        text = "\n".join(lines)
######                    except Exception as e:
######                        print(f"[ERROR] pptx extraction failed: {e}")
######                        raise
######
######                # --- XLSX/XLS ---
######                elif ext in ('.xlsx', '.xls'):
######                    try:
######                        import openpyxl
######                        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
######                        lines = []
######                        for sheet in wb.worksheets:
######                            for row in sheet.iter_rows(values_only=True):
######                                rowcells = [str(c) for c in row if c is not None and str(c).strip() != ""]
######                                if rowcells:
######                                    lines.append(" ".join(rowcells))
######                        text = "\n".join(lines)
######                    except Exception as e:
######                        print(f"[ERROR] excel extraction failed: {e}")
######                        raise
######
######                # --- Plain text / source code (include .ino) ---
######                # Add any other code extensions you want here
######                elif (
######                    (mime_type and mime_type.startswith('text'))
######                    or ext in ('.txt', '.md', '.py', '.js', '.html', '.htm', '.css',
######                               '.json', '.c', '.cpp', '.h', '.hpp', '.java', '.sh',
######                               '.bash', '.ps1', '.rb', '.go', '.rs', '.php', '.pl',
######                               '.scala', '.swift', '.kt', '.kts', '.ts', '.tsx', '.jsx',
######                               '.ino', '.ino.txt')  # .ino explicitly included
######                ):
######                    # Try utf-8 first, then chardet if available, then latin-1
######                    tried_encodings = []
######                    def _decode_try(enc):
######                        try:
######                            return file_bytes.decode(enc, errors='replace')
######                        except Exception:
######                            return None
######
######                    # quick utf-8 attempt
######                    tried_encodings.append('utf-8')
######                    text = _decode_try('utf-8')
######                    if text is None or text == "":
######                        # try chardet if present
######                        if chardet:
######                            try:
######                                detected = chardet.detect(file_bytes)
######                                enc = detected.get("encoding") or "utf-8"
######                                tried_encodings.append(enc)
######                                text = _decode_try(enc)
######                            except Exception as e:
######                                print(f"[WARN] chardet detection failed: {e}")
######                        # final fallback to latin-1
######                        if not text:
######                            tried_encodings.append('latin-1')
######                            text = _decode_try('latin-1') or ""
######                    # normalize line endings
######                    text = text.replace('\r\n', '\n').replace('\r', '\n')
######
######                # --- Image OCR ---
######                elif (mime_type and mime_type.startswith('image')) or ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'):
######                    try:
######                        from PIL import Image
######                        import pytesseract
######                        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
######                        text = pytesseract.image_to_string(img)
######                    except Exception as e:
######                        print(f"[ERROR] image OCR failed: {e}")
######                        raise
######
######                # --- Unsupported: write raw to tmp and return empty text (or optionally return path) ---
######                else:
######                    try:
######                        tmp_path = os.path.join(TMP_DIR, secure_filename(filename))
######                        with open(tmp_path, "wb") as wf:
######                            wf.write(file_bytes)
######                        print(f"[WARN] Unsupported file saved to {tmp_path}")
######                        text = ""
######                    except Exception as e:
######                        print(f"[ERROR] Could not save unsupported file: {e}")
######                        raise
######
######            except Exception as e:
######                print(f"[ERROR] extract_text_from_any failed for '{filename}': {e}")
######                traceback.print_exc()
######
######            # final cleanup & info
######            if isinstance(text, bytes):
######                try:
######                    text = text.decode("utf-8", errors="replace")
######                except Exception:
######                    text = str(text)
######
######            short = (text[:200] + "...") if len(text) > 200 else text
######            print(f"[INFO] extract_text_from_any: filename={filename!r} bytes={len(file_bytes) if file_bytes is not None else 'None'} -> chars={len(text)} preview={short!r}")
######            return text
######
######        # --- Extract text ---
######        file_text = extract_text_from_any(file_bytes, filename, mime_type)
######        file_text = re.sub(r'\s+', ' ', file_text).strip()
######        print(f"[INFO] Extracted {len(file_text)} characters from '{filename}'")
######
######        # --- Store file per-user for queries ---
######        user_files[user] = file_text
######
######        # --- Notify frontend ---
######        socketio.emit('state_update', {
######            'type': 'file_info',
######            'payload': {'message': f"Text from '{filename}' processed for embedding"},
######            'username': user
######        }, room=user)
######
######        # --- RAG EMBEDDING ---
######        try:
######            import chromadb
######
######            # ✅ FIX: use new PersistentClient if available
######            try:
######                if hasattr(chromadb, "PersistentClient"):
######                    client = chromadb.PersistentClient(path=os.path.join(BASE_DIR, "chroma_db"))
######                    print("[INFO] Using new Chroma PersistentClient")
######                else:
######                    from chromadb.config import Settings as ChromaSettings
######                    chroma_settings = ChromaSettings(
######                        chroma_db_impl="duckdb+parquet",
######                        persist_directory=os.path.join(BASE_DIR, "chroma_db")
######                    )
######                    client = chromadb.Client(chroma_settings)
######                    print("[INFO] Using legacy Chroma Client")
######            except Exception as e:
######                print(f"[WARN] Fallback to default Chroma client: {e}")
######                client = chromadb.Client()
######
######            from sentence_transformers import SentenceTransformer
######            st_model = SentenceTransformer('all-MiniLM-L6-v2')
######
######            # Chunk text for long files
######            def chunk_text(text, max_chars=1500):
######                for i in range(0, len(text), max_chars):
######                    yield text[i:i+max_chars]
######
######            safe_name = secure_filename(filename)
######            doc_id = f"{user}_{safe_name}_{int(time.time())}"
######
######            # ✅ version-safe collection get/create
######            if hasattr(client, "get_or_create_collection"):
######                collection = client.get_or_create_collection(name="docs")
######            elif hasattr(client, "create_collection"):
######                collection = client.create_collection(name="docs")
######            else:
######                collection = client.get_collection("docs")
######
######            for i, chunk in enumerate(chunk_text(file_text)):
######                vect = st_model.encode(chunk)
######                embedding = vect.tolist() if hasattr(vect, "tolist") else list(vect)
######                try:
######                    collection.add(
######                        ids=[f"{doc_id}_chunk{i}"],
######                        embeddings=[embedding],
######                        documents=[chunk]
######                    )
######                except Exception:
######                    # Try alternate add signature (for newer Chroma versions)
######                    collection.add(
######                        documents=[chunk],
######                        metadatas=[{"source": filename}],
######                        ids=[f"{doc_id}_chunk{i}"],
######                        embeddings=[embedding]
######                    )
######
######            # ✅ persist safely
######            try:
######                if hasattr(client, "persist"):
######                    client.persist()
######                elif hasattr(collection, "persist"):
######                    collection.persist()
######            except Exception as e:
######                print(f"[WARN] Persist skipped: {e}")
######
######            print(f"[INFO] Embedded {len(file_text)} chars into Chroma for RAG")
######
######            # Log event
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": f"[file_upload] {filename}",
######                "response": "[stored in RAG]",
######                "models": {"embed_model": "sentence-transformers"}
######            }
######            log_user_query(user, entry)
######
######            # Notify frontend success
######            socketio.emit('state_update', {
######                'type': 'file_info',
######                'payload': {
######                    'message': f"Text '{filename}' processed for RAG",
######                    'doc_id': doc_id,
######                    'length': len(file_text),
######                    'rag_available': True,
######                    'embed_source': "sentence-transformers"
######                },
######                'username': user
######            }, room=user)
######
######        except Exception as e:
######            print(f"[ERROR] RAG embedding failed: {e}")
######            socketio.emit('state_update', {
######                'type': 'file_error',
######                'payload': f'RAG processing failed: {e}',
######                'username': user
######            }, room=user)
######
######        return
######
######
######    # --- QUERY ---
######    if t == 'query':
######        text = payload if isinstance(payload, str) else payload.get('query', '')
######        print(f"[EVENT] query from user='{user}' text='{text}'")
######
######        if not text:
######            print("[WARN] Empty query text.")
######            return
######
######        user_files = globals().get("user_files", {})
######        if any(x in text.lower() for x in ["this file", "this document", "this code", "this text"]) and user in user_files:
######            try:
######                from rag_embeddings import file_embeddings
######                fe = file_embeddings()
######                fe.build_store(user_files[user])
######                context_chunks = fe.retrieve(query=text, top_k=20)
######                context = "\n".join(f"- {chunk.strip()}" for chunk in context_chunks) if context_chunks else "No relevant items found."
######
######                final_prompt = f"{text}\n\nContext:\n{context}"
######                print(f"[RAG QUERY] Final prompt:\n{final_prompt[:400]}...")
######
######            except Exception as e:
######                print(f"[RAG ERROR QUERY] {e}")
######                final_prompt = text
######
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": final_prompt,
######                "response": None,
######                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
######            }
######            log_user_query(user, entry)
######            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
######            return
######
######        else:
######            entry = {
######                "username": user,
######                "timestamp": datetime.now().isoformat(),
######                "query": text,
######                "response": None,
######                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
######            }
######            log_user_query(user, entry)
######            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
######            return
######
######
####### HTTP endpoints
######@app.route('/main_responses/<username>', methods=['GET'])
######def main_responses(username):
######    path = os.path.join(TMP_DIR, f"{username}_main.json")
######    if not os.path.exists(path):
######        return jsonify([])
######    with open(path, 'r', encoding='utf-8') as f:
######        try:
######            return jsonify(json.load(f))
######        except json.JSONDecodeError:
######            return jsonify([])
######
######@app.route('/user_logs/<username>', methods=["GET"])
######def get_user_logs(username):
######    logs = read_logs(username)
######    print(f"[DEBUG] Loaded logs for {username}: {len(logs)} entries")
######    return jsonify(logs)
######
######@app.post("/upload_audio")
######def upload_audio():
######    os.makedirs(TMP_DIR, exist_ok=True)
######    webm_path = os.path.join(TMP_DIR, "temp_audio.webm")
######    wav_path = os.path.join(TMP_DIR, "temp_audio.wav")
######
######    if 'audio' not in request.files:
######        return jsonify({"error": "No audio file"}), 400
######    request.files["audio"].save(webm_path)
######
######    # convert with ffmpeg if available
######    if os.system(f'ffmpeg -y -i "{webm_path}" -ar 16000 -ac 1 -f wav "{wav_path}"') != 0:
######        return jsonify({"error": "ffmpeg failed"}), 500
######
######    try:
######        text = whisper_model.transcribe(wav_path).get("text", "").strip()
######        cleaned = webui_listener.listen_text(text)
######    except Exception:
######        cleaned = ""
######
######    for p in (webm_path, wav_path):
######        try:
######            os.remove(p)
######        except:
######            pass
######
######    return jsonify({"transcript": cleaned})
######
######@app.post('/upload_file')
######def upload_file_http():
######    # HTTP multipart upload alternative to socket-based upload
######    if 'file' not in request.files:
######        return jsonify({"error": "No file provided"}), 400
######    file = request.files['file']
######    username = request.form.get('username', 'Tjaart')
######    filename = secure_filename(file.filename)
######    tmp_path = os.path.join(TMP_DIR, filename)
######    file.save(tmp_path)
######
######    # Try to detect type and if text, process into Chroma
######    mime_type, _ = mimetypes.guess_type(filename)
######
######    if mime_type and mime_type.startswith('image'):
######        print(f"[HTTP FILE] Image uploaded: {filename} by {username}")
######        socketio.emit('state_update', {'type': 'file_info', 'payload': f"Image '{filename}' received", 'username': username}, room=username)
######        return jsonify({"status": "ok", "type": "image"})
######
######    if mime_type and (mime_type.startswith('text') or filename.lower().endswith(('.txt', '.md', '.py'))):
######        print(f"[HTTP FILE] Text uploaded: {filename} by {username}")
######        try:
######            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
######                file_text = f.read()
######        except Exception as e:
######            return jsonify({"error": f"Could not read file: {e}"}), 500
######
######        # Same RAG logic as socket handler
######        try:
######            client = None
######            if chromadb is not None:
######                try:
######                    from chromadb.config import Settings as ChromaSettings
######                    chroma_settings = ChromaSettings(
######                        chroma_db_impl="duckdb+parquet",
######                        persist_directory=os.path.join(BASE_DIR, "chroma_db")
######                    )
######                    client = chromadb.Client(chroma_settings)
######                except Exception:
######                    try:
######                        client = chromadb.Client()
######                    except Exception:
######                        client = None
######
######            embedding = None
######            emb_source = None
######            try:
######                if ollama is not None:
########                    emb_resp = ollama.embeddings(model="mxbai-embed-large", prompt=file_text)
######                    emb_resp = ollama.embeddings(model="mxbai-embed-large")
######                    embedding = emb_resp.get('embedding')
######                    emb_source = 'ollama'
######            except Exception:
######                try:
######                    from sentence_transformers import SentenceTransformer
######                    st_model = SentenceTransformer('all-MiniLM-L6-v2')
######                    vect = st_model.encode(file_text)
######                    embedding = vect.tolist() if hasattr(vect, 'tolist') else list(vect)
######                    emb_source = 'sentence-transformers'
######                except Exception:
######                    embedding = None
######                    emb_source = None
######
######            doc_id = f"{username}_{filename}_{int(time.time())}"
######            if embedding is not None and client is not None:
######                try:
######                    collection = client.get_or_create_collection(name="docs")
######                    collection.add(ids=[doc_id], embeddings=[embedding], documents=[file_text])
######                    try:
######                        client.persist()
######                    except Exception:
######                        pass
######                except Exception:
######                    client = None
######
######            if embedding is None or client is None:
######                fallback_index = os.path.join(TMP_DIR, "local_rag_index.json")
######                try:
######                    if os.path.exists(fallback_index):
######                        with open(fallback_index, 'r', encoding='utf-8') as fh:
######                            idx = json.load(fh)
######                    else:
######                        idx = []
######                except Exception:
######                    idx = []
######                idx_entry = {
######                    'doc_id': doc_id,
######                    'username': username,
######                    'filename': filename,
######                    'path': tmp_path,
######                    'length': len(file_text),
######                    'embedding_source': emb_source,
######                    'timestamp': datetime.now().isoformat()
######                }
######                idx.append(idx_entry)
######                try:
######                    with open(fallback_index, 'w', encoding='utf-8') as fh:
######                        json.dump(idx, fh, ensure_ascii=False, indent=2)
######                except Exception:
######                    pass
######
######            log_user_query(username, {"username": username, "timestamp": datetime.now().isoformat(), "query": f"[file_upload] {filename}", "response": "[stored in RAG]" if embedding is not None else "[saved - RAG unavailable]", "models": {"embed_model": emb_source or "none"}})
######
######    ##        prompt=f"Using this data: {data}. Respond to this prompt: {prompt}",
######            full_emb_prompt = f"Using this data: {response_text}. Respond to this prompt: {text}"
######            socketio.emit('state_update', {
######                'type': 'query',
######    ##            'payload': {'query': text, 'response': response_text},
######                'payload': full_emb_prompt,
######                'username': user
######            })
######
########            socketio.emit('state_update', {'type': 'file_info', 'payload': {'message': f"Text '{filename}' processed for RAG" if embedding is not None else f"Text '{filename}' saved (RAG unavailable)", 'doc_id': doc_id}, 'username': username}, room=username)
######            return jsonify({"status": "ok", "type": "text", "doc_id": doc_id})
######        except Exception as e:
######            print(f"[ERROR] HTTP RAG failed: {e}")
######            return jsonify({"error": str(e)}), 500
######
######    return jsonify({"status": "ok", "type": "unknown"})
######
######
######@app.route('/edge_tts', methods=['POST'])
######def edge_tts_endpoint():
######    """
######    Synthesize text with edge-tts and return MP3 bytes.
######    Frontend posts JSON: { text, voice?, style? }.
######    Returns: audio/mpeg binary (200) or JSON error (non-200).
######    """
######    data = request.get_json(silent=True) or {}
######    text = data.get('text') or data.get('response') or ""
######    voice = data.get('voice', 'en-US-GuyNeural')
######    style = data.get('style', None)  # currently unused in this simple example
######
######    if not text:
######        return jsonify({"error": "No 'text' provided"}), 400
######
######    # Try importing edge-tts
######    try:
######        import edge_tts
######    except Exception as e:
######        return jsonify({"error": f"edge-tts not installed or import failed: {e}"}), 500
######
######    # create a temporary file for the mp3
######    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
######    tmp_path = tmp.name
######    tmp.close()
######
######    async def synth_to_file(txt, voice_name):
######        # Use simple text; if you want SSML or style-based SSML, build it here.
######        comm = edge_tts.Communicate(txt, voice_name)
######        # save will write mp3 to path
######        await comm.save(tmp_path)
######
######    try:
######        # Run the async synth (blocks until complete)
######        asyncio.run(synth_to_file(text, voice))
######    except Exception as e:
######        # cleanup
######        try:
######            os.unlink(tmp_path)
######        except Exception:
######            pass
######        return jsonify({"error": f"edge-tts synthesis failed: {e}"}), 500
######
######    # Return the mp3 file bytes
######    try:
######        # send_file will set appropriate headers
######        resp = send_file(tmp_path, mimetype='audio/mpeg', as_attachment=False)
######        # Allow cross-origin if needed (your SocketIO allowed * already)
######        resp.headers['Access-Control-Allow-Origin'] = '*'
######        return resp
######    finally:
######        # cleanup file after response has been scheduled (Flask may still be reading it,
######        # but removing the file here is okay on many OSes; if you prefer, delete in a worker)
######        try:
######            os.unlink(tmp_path)
######        except Exception:
######            pass
######
######@app.route('/login', methods=["POST"])
######def login_user():
######    data = request.json
######    success, msg = do_login(data["username"], data["password"])
######    return jsonify({"success": success, "message": msg})
######
######@app.route('/signup', methods=["POST"])
######def signup_user():
######    data = request.json
######    success, msg = do_signup(data["username"], data["password"])
######    return jsonify({"success": success, "message": msg})
######
######@app.route('/users', methods=["GET"]) 
######def list_users():
######    return jsonify(get_users())
######
######@app.route('/delete_user', methods=["POST"]) 
######def remove_user():
######    data = request.json
######    success = delete_user(data["username"])
######    return jsonify({"success": success})
######
######@app.route('/', defaults={'path': ''})
######@app.route('/<path:path>')
######def serve(path):
######    if path != "" and os.path.exists(os.path.join(app.static_folder, path)):
######        return send_from_directory(app.static_folder, path)
######    return send_from_directory(app.static_folder, 'index.html')
######
######if __name__ == '__main__':
######    context = ('cert.pem', 'key.pem') if os.path.exists('cert.pem') and os.path.exists('key.pem') else None
######    socketio.run(app, host='0.0.0.0', port=5001, ssl_context=context)
######



##      # LASTEST WORKING 
##
##import ollama
##import re
##import os
##import json
##import time
##import base64
##import mimetypes
##import tempfile
##from datetime import datetime
##from werkzeug.utils import secure_filename
##
##import whisper
##from flask import Flask, request, jsonify, send_from_directory
##from flask_socketio import SocketIO, emit, join_room, leave_room
##
### Optional libs (may not be installed in every env)
##try:
##    import chromadb
##except Exception:
##    chromadb = None
##try:
##    import ollama
##except Exception:
##    ollama = None
##
##from listenWEBUI import WEBUIListenModule
##from auth_manager import login as do_login, signup as do_signup, get_users, delete_user
##from query_logger import log_user_query, log_user_response, read_logs
##
##import chardet
##from pdfminer.high_level import extract_text
##from pdf2image import convert_from_path
##import pytesseract
##import base64
##import io
##import PyPDF2
##from PyPDF2 import PdfFileReader
##import pdfplumber
##from langchain_chroma import Chroma
##
##import asyncio
##import tempfile
##from flask import send_file   # modify existing import line to include send_file if needed
##
##
##
##BASE_DIR = os.path.dirname(os.path.abspath(__file__))
##TMP_DIR = os.path.join(BASE_DIR, 'tmp')
##os.makedirs(TMP_DIR, exist_ok=True)
##
### Initialize listener and Whisper
##webui_listener = WEBUIListenModule()
##whisper_model = whisper.load_model("base.en")
##
##app = Flask(__name__, static_folder='./webui-src/dist', static_url_path='')
##socketio = SocketIO(app, cors_allowed_origins='*', async_mode='threading')
##
##
##chromadb.api.client.SharedSystemClient.telemetry_enabled = False
##
##
### Module-level state
##state = {
##    'logs': [],
##    'queries': [],
##    'resp': [],
##    'settings': {
##        'use_whisper': False,
##        'use_vosk': True,
##        'enter_submits': True
##    }
##}
##
### Session tracking
##session_users = {}  # Maps SID -> username
##
##my_file_content = ""
##
##collection = None   # Chroma collection shared across events
##os.environ["CHROMA_TELEMETRY_ENABLED"] = "false"
##
##from langchain.text_splitter import RecursiveCharacterTextSplitter
##import PyPDF2
##
##
##model_selected = ""
##thinkbot_model = ""
##rag_clear = False
##detected_model = ""
##
### ======= Query helper (can be called from elsewhere in backend) =======
##def query_with_retrieval(prompt: str, n_results: int = 1):
##    """Retrieve top chunks from Chroma and run generation with Ollama."""
##    try:
##        if ollama is None or chromadb is None:
##            raise RuntimeError("Ollama or Chroma not available")
##
##        q_emb = ollama.embeddings(model="mxbai-embed-large", prompt=prompt)
##        q_vector = q_emb.get("embedding") if isinstance(q_emb, dict) else q_emb
##        client = chromadb.Client()
##        collection = client.get_or_create_collection(name="docs")
##        res = collection.query(query_embeddings=[q_vector], n_results=n_results)
##        top_doc = res.get("documents", [[]])[0][0] if res.get("documents") else ""
##    except Exception as e:
##        print(f"[WARN] retrieval failed: {e}")
##        top_doc = ""
##
##    gen_prompt = f"Using this data: {top_doc}\n\nRespond to: {prompt}"
##    try:
##        out = ollama.generate(model="tinyllama", prompt=gen_prompt, stream=False)
##        text = out.get("response") or out.get("text") or str(out)
##    except Exception as e:
##        text = f"[generation failed: {e}]"
##    return text
##
##def extract_pdf_text_with_ocr(pdf_path):
##    try:
##        text = extract_text(pdf_path)
##        if len(text.strip()) > 10:  # if decent text extracted, return it
##            return text
##        else:
##            # Fallback to OCR if text is too short or empty
##            print("[INFO] PDF text extraction empty or too short, using OCR fallback.")
##            pages = convert_from_path(pdf_path)
##            ocr_text = ""
##            for page in pages:
##                ocr_text += pytesseract.image_to_string(page)
##            return ocr_text
##    except Exception as e:
##        print(f"[ERROR] PDF extraction + OCR failed: {e}")
##        return ""
##
##
### Utility: safe save
##def _save_tmp(filename, data_bytes):
##    safe = secure_filename(filename)
##    path = os.path.join(TMP_DIR, safe)
##    with open(path, 'wb') as f:
##        f.write(data_bytes)
##    return path
##
### Socket handlers
##@socketio.on('connect')
##def connect():
##    emit('full_state', state)
##
##@socketio.on('disconnect')
##def disconnect():
##    sid = request.sid
##    user = session_users.pop(sid, None)
##    if user:
##        leave_room(user)
##        print(f"[DISCONNECT] {sid} left room {user}")
##
##
### --- helper: singleton chroma client creator (safe) ---
##
##_CHROMA_CLIENT = None
##
##def get_chroma_client(persist_directory=None, chroma_db_impl="duckdb+parquet"):
##    global _CHROMA_CLIENT
##    if _CHROMA_CLIENT is not None:
##        return _CHROMA_CLIENT
##    try:
##        import chromadb
##        from chromadb.config import Settings as ChromaSettings
##    except Exception as e:
##        print(f"[RAG] chromadb import failed: {e}")
##        return None
##
##    try:
##        settings = ChromaSettings(chroma_db_impl=chroma_db_impl,
##                                  persist_directory=persist_directory)
##
##        client = None
##        try:
####            from chromadb.config import Settings as ChromaSettings
##            chroma_settings = ChromaSettings(
##                chroma_db_impl="duckdb+parquet",
##                persist_directory=os.path.join(BASE_DIR, "chroma_db")
##            )
##            client = chromadb.Client(chroma_settings)
##        except Exception:
##            # fallback to default client constructor; may still fail, handle below
##            try:
##                client = chromadb.Client()
##            except Exception as e:
##                client = None
##                print(f"[RAG] chromadb.Client fallback failed: {e}")
##
##
##            if client is None:
##                client = get_chroma_client(persist_directory=CHROMA_DIR)
##
##        return client
##
##
##
##
####        _CHROMA_CLIENT = chromadb.Client(settings=settings)
####        print(f"[RAG] Chroma client created with persist_directory={persist_directory}")
####        return _CHROMA_CLIENT
##    except Exception as e:
##        print(f"[RAG] Embedded Chroma init failed: {e}; trying fallback chromadb.Client()")
##        try:
##            _CHROMA_CLIENT = chromadb.Client()
##            print("[RAG] Chroma client created via fallback chromadb.Client()")
##            return _CHROMA_CLIENT
##        except Exception as e2:
##            print(f"[RAG] chromadb.Client() fallback failed: {e2}")
##            _CHROMA_CLIENT = None
##            return None
##
##_MODEL_STR_PATTERN = re.compile(r"^[A-Za-z0-9_\-]+(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)?$")
##
##def extract_model(payload) -> str:
##    """
##    Robustly extract and clean a model name. Always returns a string ("" if none found).
##
##    Rules:
##      - If payload is a dict and contains thinkbot_model/model/model_used -> use it.
##      - Else if payload is a dict and has 'payload' whose value is a dict -> descend and try again.
##      - Else if payload is a plain string and *looks like a model name* (pattern) -> accept it.
##      - Otherwise return "".
##
##    Example results:
##      {'payload': {'thinkbot_model': 'moondream'}} -> "moondream"
##      {'thinkbot_model': 'deepseek-r1-1.5b'}    -> "deepseek-r1"
##      "deepseek-r1"                             -> "deepseek-r1"
##      {'type': 'log', 'payload': 'I see, a man ...'} -> ""
##    """
##    model_raw = ""
##
##    # If dict, prefer direct model keys at top-level
##    if isinstance(payload, dict):
##        # Direct model keys first
##        if "detected_model" in payload or "model" in payload or "model_used" in payload:
##            model_raw = payload.get("detected_model") or payload.get("model") or payload.get("model_used") or ""
##        else:
##            # Only descend into 'payload' if it's a dict (avoid descending into plain log strings)
##            inner = payload.get("payload")
##            if isinstance(inner, dict):
##                model_raw = inner.get("detected_model") or inner.get("model") or inner.get("model_used") or ""
##            else:
##                # No model found in this event
##                return ""
##    else:
##        # payload not a dict; if it's a plain string that *looks like a model* accept it
##        if isinstance(payload, str) and _MODEL_STR_PATTERN.match(payload.strip()):
##            model_raw = payload.strip()
##        else:
##            return ""
##
##    model_raw = (model_raw or "").strip()
##    if not model_raw:
##        return ""
##
##    # Clean suffixes like ":1.5b", "-1.5b", ":8b", "-8b", etc.
##    model_clean = re.sub(r"(?:(?:[:\-])\d+(?:\.\d+)?[A-Za-z]*)$", "", model_raw)
##
##    return model_clean
##
##
##@socketio.on('gui_event')
##def handle_event(data):
##    global client, collection, detected_model, thinkbot_model, model_selected, rag_clear  # Declare globals at the very top
##    sid = request.sid
##    print(f"[DEBUG] Received gui_event from sid={sid}: {data!r}")
##
##    user = data.get('username') if isinstance(data, dict) else None
##    t = data.get('type') if isinstance(data, dict) else None
##    payload = data.get('payload') if isinstance(data, dict) else None
##
##    event = data
##    print("[EVENT] Received event 'new event' : ", str(event))
##
##    model_used = extract_model(event)
##    print(f"[EVENT] Received gui_event 'model_used' : {model_used}")
##
##    if not user:
##        user = 'Tjaart'  # fallback
##    # Track user per session
##    session_users[sid] = user
##    join_room(user)
##    global my_file_content
##    
##    # --- LOGIN ---
##    if t == 'login':
##        print(f"[LOGIN] User: {user}")
##        socketio.emit('state_update', {
##            'type': 'login_ack',
##            'username': user
##        }, room=user)
##        return
##
##    # --- LOG ---
##    if t == 'log':
##        # Always produce a structured entry so frontend can read models reliably
##        if isinstance(payload, str):
##            # Keep old logging side-effect
##            log_user_response(user, payload)
##
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": "",
##                "response": payload,
##                "models": {"thinkbot": model_used or ""}
##            }
##
##            socketio.emit('state_update', {
##                'type': 'log',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##        elif isinstance(payload, dict):
##            # Ensure models is a dict and add thinkbot
##            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
##            models = dict(models)  # copy-safe
##            models["thinkbot"] = model_used or ""
##            query = payload.get("query", "")
##
##            if query:
##                guery = query.replace(":","")
##                print(f"New Query supplied : {query}")
##            else:
##                print("There is no query supplied...")
##
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": guery,
##                "response": payload.get("response", ""),
##                "models": models
##            }
##
##            log_user_query(user, entry)
##
##            socketio.emit('state_update', {
##                'type': 'log',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##            socketio.emit('state_update', {
##                'type': 'query',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##        return
##
##    # --- RESP ---
##    if t == 'resp':
##        if isinstance(payload, str):
##            log_user_response(user, payload)
##
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": "",
##                "response": payload,
##                "models": {"thinkbot": model_used or ""}
##            }
##
##            socketio.emit('state_update', {
##                'type': 'log',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##        elif isinstance(payload, dict):
##            # Preserve any existing models and inject thinkbot
##            models = payload.get("models") if isinstance(payload.get("models"), dict) else {}
##            models = dict(models)
##            models["thinkbot"] = model_used or ""
##            query = payload.get("query", "")
##            
##            if query:
##                query = query.replace(":","")
##                print(f"New Query supplied : {query}")
##            else:
##                print("No Query supplied...")
##                 
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": query,
##                "response": payload.get("response", ""),
##                "models": models
##            }
##
##            log_user_query(user, entry)
##
##            socketio.emit('state_update', {
##                'type': 'log',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##            socketio.emit('state_update', {
##                'type': 'query',
##                'payload': entry,
##                'username': user
##            }, room=user)
##
##            print(f" \n thinkbot_model : {model_used} \n")
##
##        return
##
##    # --- SETTING ---
##    if t == 'setting':
##        socketio.emit('state_update', {
##            'type': 'setting',
##            'payload': payload,
##            'username': user
##        }, room=user)
##        return
##
##
##    if t in ('setting', 'detected_model'):
##        data_payload = data.get('payload')
##        user = data.get('username') or 'default_user'
####        new_detected_model = data.get('detected_model')
##
##        # Try to extract the thinkbot model name from different payload shapes:
##        model_selected = None
##        
##        if isinstance(data_payload, dict):
##            # common: payload = { 'detected_model': 'tinyllama' }
##            if 'detected_model' in data_payload:
##                model_selected = data_payload['detected_model']
##            else:
##                # fallback: payload may be { '<some_model_key>': '<name>' }
##                # find any key that ends with '_model' and pick its value
##                model_keys = [k for k in data_payload.keys() if k.endswith('_model')]
##                if len(model_keys) == 1:
##                    model_selected = data_payload[model_keys[0]]
##        elif isinstance(data_payload, str):
##            # sometimes frontend sends a raw string payload (e.g. from 'think_model' emit)
##            model_selected = data_payload
##
##        # If we found a model, assign and notify
##        if model_selected:
##            # ensure we set the module/global variable (declare `global detected_model` above if needed)
##            try:
##                global detected_model
##            except NameError:
##                # if detected_model doesn't exist yet, create it
##                detected_model = None
##
##            detected_model = model_selected  # correct assignment (not ==)
##            print(f"\n The detected_model for RAG is '{detected_model} \n'")
##            print(f"\n The new detected_model SETTINGS is now '{detected_model}' \n")
##
##            # Emit a consistent state_update so clients can sync
##            socketio.emit('state_update', {
##                'type': 'setting',
##                'payload': {'detected_model': detected_model},
##                'username': user
##            }, room=user)
##        else:
##            # No model found — still forward the raw payload for other settings
##            print(f"[setting] received payload but no model detected: {data_payload}")
##            socketio.emit('state_update', {
##                'type': 'setting',
##                'payload': data_payload,
##                'username': user
##            }, room=user)
##
##        return
##
##
##
####    if t in 'detected_model':
####
####        payload = data.get('payload')
####        user = data.get('username') or 'default_user'
####        new_detected_model = data.get('detected_model')
####
####        print(f"The model for new_detected_model 'model_selected' is {new_detected_model}")
####        return
##
##
##    # --- SETTING ---
##    if t == 'think_model':
##
##        payload = data.get('payload')
##        user = data.get('username') or 'default_user'
##        model_selected = data.get('think_model')
##
##        print(f"The model for RAG 'model_selected' is {model_selected}")
##        return
##
##    # --- RESET RAG ---
##    if t == 'reset_rag':
##        import shutil
##        # ensure globals we touch are declared
##        global _CHROMA_CLIENT, _CHROMA_CLIENT_SETTINGS, collection, my_file_content, RAG_STORE
##        user = data.get('username') if isinstance(data, dict) else user
##        user = user or getpass.getuser()
##        print(f"[RAG] reset_rag requested by {user}")
##
##        # canonical persist dir used everywhere in this app
##        CHROMA_DIR = os.path.join(BASE_DIR, "chroma_store")
##        print(f"[RAG] CHROMA_DIR for RAG is {CHROMA_DIR}")
##
##        # Ensure client exists in local scope (avoid NameError)
##        client = None
##
##        # Try to get (or lazily create) a client via the cached helper
##        try:
##            client = get_chroma_client(persist_directory=CHROMA_DIR)
##        except Exception as e:
##            print(f"[RAG] get_chroma_client() raised during reset attempt: {e}")
##            client = None
##
##        try:
##            # If chromadb isn't present, bail with an informative error
##            if chromadb is None:
##                raise RuntimeError("chromadb not available in this environment")
##
##            # If we didn't manage to get a client, try again (lazy)
##            if client is None:
##                client = get_chroma_client(persist_directory=CHROMA_DIR)
##
##            # 1) Preferred API: delete the collection if client exposes it
##            deleted = False
##            if client is not None and hasattr(client, "delete_collection"):
##                try:
##                    client.delete_collection("docs")
##                    print("[RAG] client.delete_collection('docs') succeeded")
##                    deleted = True
##                except Exception as exc:
##                    print(f"[RAG] client.delete_collection failed: {exc!r}")
##
##            # 2) If not deleted above: attempt to obtain the collection and clear it defensively
##            if not deleted and client is not None:
##                try:
##                    # use get_or_create_collection for max compatibility
##                    if hasattr(client, "get_or_create_collection"):
##                        coll = client.get_or_create_collection(name="docs")
##                    else:
##                        coll = client.get_collection("docs") if hasattr(client, "get_collection") else None
##                        if coll is None and hasattr(client, "create_collection"):
##                            coll = client.create_collection("docs")
##                    # try a no-arg delete()
##                    try:
##                        if coll is not None and hasattr(coll, "delete"):
##                            coll.delete()  # many versions will clear all docs
##                            print("[RAG] collection.delete() called successfully (if supported)")
##                            deleted = True
##                    except Exception as exc:
##                        print(f"[RAG] collection.delete() failed: {exc!r}")
##
##                    # If still not deleted, try to fetch all ids and delete by ids
##                    if not deleted and coll is not None:
##                        try:
##                            # attempt to get all ids via coll.get(include=['ids'])
##                            ids = []
##                            try:
##                                res = coll.get(include=['ids'])
##                                ids_list = res.get('ids', [[]]) if isinstance(res, dict) else []
##                                # normalize shape: often it's [ [id1, id2, ...] ]
##                                if isinstance(ids_list, list) and len(ids_list) > 0 and isinstance(ids_list[0], list):
##                                    ids = ids_list[0]
##                                elif isinstance(ids_list, list):
##                                    ids = ids_list
##                            except Exception:
##                                ids = []
##
##                            if ids:
##                                try:
##                                    coll.delete(ids=ids)
##                                    print(f"[RAG] collection.delete(ids=...) succeeded ({len(ids)} ids)")
##                                    deleted = True
##                                except Exception as exc:
##                                    print(f"[RAG] collection.delete(ids=...) failed: {exc!r}")
##                        except Exception as exc:
##                            print(f"[RAG] Attempt to list+delete ids failed: {exc!r}")
##
##                except Exception as exc_coll:
##                    print(f"[RAG] failure while attempting to clear collection via client: {exc_coll!r}")
##
##            # 3) If none of the above deleted the store, fallback to deleting the persist directory
##            if not deleted:
##                try:
##                    if os.path.exists(CHROMA_DIR):
##                        shutil.rmtree(CHROMA_DIR)
##                        os.makedirs(CHROMA_DIR, exist_ok=True)
##                        print(f"[RAG] Persist directory {CHROMA_DIR} removed and recreated (last-resort cleanup)")
##                    else:
##                        print(f"[RAG] Persist directory {CHROMA_DIR} did not exist; nothing to remove")
##                except Exception as exc_sh:
##                    raise RuntimeError(f"Failed to remove persist directory fallback: {exc_sh}")
##
##            # Clear cached client and local references so process-level state is reset
##            _CHROMA_CLIENT = None
##            _CHROMA_CLIENT_SETTINGS = None
##            client = None
##            collection = None
##
##            # Clear in-memory RAG references (so UI can't still use old chunks)
##            try:
##                if isinstance(RAG_STORE, dict) and user in RAG_STORE:
##                    del RAG_STORE[user]
##            except Exception:
##                RAG_STORE = {}
##
##            # Optionally clear global last-uploaded file (if you want)
##            try:
##                my_file_content = ""
##            except Exception:
##                pass
##
##            # Recreate fresh client & empty collection so system is immediately usable
##            client = get_chroma_client(persist_directory=CHROMA_DIR)
##            if client is not None:
##                try:
##                    collection = client.get_or_create_collection(name="docs")
##                    # defensive: try removing docs again if collection.delete exists
##                    try:
##                        if hasattr(collection, "delete"):
##                            collection.delete()
##                    except Exception:
##                        pass
##                except Exception as e_create:
##                    print(f"[RAG] Warning: could not create/clear collection after reset: {e_create}")
##
##            print("[RAG] Reset completed successfully")
##            socketio.emit('state_update', {
##                'type': 'reset_rag_ack',
##                'payload': 'RAG index cleared',
##                'username': user
##            }, room=user)
##
##        except Exception as e:
##            print(f"[RAG] Reset failed: {e}")
##            socketio.emit('state_update', {
##                'type': 'reset_rag_error',
##                'payload': f'RAG reset failed: {e}',
##                'username': user
##            }, room=user)
##
##        return
##
##    # --- LOAD FILE ---
##
##    if t == 'load_file':
##        my_file_content = ""
##
##        # Validate payload
##        if not isinstance(payload, dict):
##            print(f"[ERROR] load_file called with invalid payload: {payload!r}")
##            socketio.emit('state_update', {
##                'type': 'file_error',
##                'payload': "Invalid or missing payload for load_file",
##                'username': user
##            }, room=user)
##            return
##
##        filename = payload.get('filename')
##        filedata_b64 = payload.get('filedata')
##        mime_hint = payload.get('mime') or payload.get('type')
##
##        print(f"[DEBUG] load_file payload keys: {list(payload.keys())}")
##
##        if not filename or filedata_b64 is None:
##            print(f"[ERROR] load_file missing filename or filedata. filename={filename!r}, filedata_present={filedata_b64 is not None}")
##            socketio.emit('state_update', {
##                'type': 'file_error',
##                'payload': {
##                    'message': 'Missing filename or filedata',
##                    'received_keys': list(payload.keys())
##                },
##                'username': user
##            }, room=user)
##            return
##
##        try:
##            file_bytes = base64.b64decode(filedata_b64)
##        except Exception as e:
##            print(f"[ERROR] base64 decode failed: {e}; payload_snippet={str(filedata_b64)[:80]}")
##            socketio.emit('state_update', {'type': 'file_error', 'payload': f'base64 decode failed: {e}', 'username': user}, room=user)
##            return
##
##        # Detect MIME type
##        mime_type, _ = mimetypes.guess_type(filename)
##        if not mime_type and mime_hint:
##            mime_type = mime_hint
##
##        # Handle PDF files directly in memory
##        if mime_type == 'application/pdf' or filename.lower().endswith('.pdf'):
##            try:
##                pdf_stream = io.BytesIO(file_bytes)
##                reader = PyPDF2.PdfReader(pdf_stream)
##                extracted_text = ""
##                for page in reader.pages:
##                    page_text = page.extract_text()
##                    if page_text:
##                        extracted_text += page_text + "\n"
##                my_file_content = extracted_text.encode('utf-8', errors='replace').decode('utf-8')
##                print(f"[INFO] Extracted text from PDF ({len(my_file_content)} chars)")
##            except Exception as e:
##                print(f"[ERROR] Failed to extract text from PDF: {e}")
##                my_file_content = ""
##
##        # Handle text files
##        elif mime_type and (mime_type.startswith('text') or filename.lower().endswith(('.txt', '.md', '.py'))):
##            try:
##                result = chardet.detect(file_bytes)
##                detected_encoding = result['encoding'] or 'utf-8'
##                decoded_content = file_bytes.decode(detected_encoding, errors='replace')
##                my_file_content = decoded_content.encode('utf-8').decode('utf-8')
##                print(f"[INFO] Detected encoding: {detected_encoding}")
##            except Exception as e:
##                print(f"[ERROR] Failed to decode text file: {e}")
##
##                # Path to your .py file
##                file_path = "d:/Python_Env/New_Virtual_Env/Alfred_Offline_New_GUI/2025_08_17_WEBUI_RAG_New/New_V2_Home_Head_Movement_Smoothing/modules/backend.py"
##
##                # Open and read the file
##                with open(file_path, "r", encoding="utf-8") as f:
##                    document_text = f.read()
##
##                print(f"document_text: {document_text[:500]}...")  # print first 500 chars for brevity
##
##                # Join text (here it's already one string, so this does nothing)
##                My_Joined_Text = document_text
##                print(f"My_Joined_Text: {My_Joined_Text[:500]}...")  # again first 500 chars
##
##                # Put it into a list
##                documents = [My_Joined_Text]
##                print(f"documents: {documents[:1]}")  # show only first element
##
##                my_file_content =  documents
##                print(f"my_file_content: {my_file_content}")  # show only first element
##
##
##        else:
##            # For unknown types, save the file for later or ignore text extraction
##            safe_name = secure_filename(filename)
##            tmp_path = os.path.join(TMP_DIR, safe_name)
##            try:
##                with open(tmp_path, 'wb') as f:
##                    f.write(file_bytes)
##                print(f"[INFO] Saved unknown file type to {tmp_path}")
##            except Exception as e:
##                print(f"[ERROR] Failed to save unknown file type: {e}")
##
##            my_file_content = ""
##
##        # Use extracted text or fallback decode
##        if my_file_content:
##            file_text = my_file_content
##        else:
##            try:
##                file_text = file_bytes.decode('utf-8', errors='ignore')
##            except Exception:
##                file_text = file_bytes.decode('latin-1', errors='ignore')
##
##        print(f"file_text is now decoded is before {file_text}")
##        
##        # (Example: log and emit success)
##        doc_id = f"{user}_{filename}_{int(time.time())}"
##        print(f"[INFO] Ready to embed document id={doc_id} with length={len(file_text)}")
##        socketio.emit('state_update', {
##            'type': 'file_info',
##            'payload': {
##                'message': f"Text from '{filename}' processed for embedding",
##            },
##            'username': user
##        }, room=user)
##
##
##        # RAG embedding & storage (same as your existing logic)
##        try:
##            # Initialize Chroma
##            try:
##                from chromadb.config import Settings as ChromaSettings
##                chroma_settings = ChromaSettings(
##                    chroma_db_impl="duckdb+parquet",
##                    persist_directory=os.path.join(BASE_DIR, "chroma_db")
##                )
##                client = chromadb.Client(chroma_settings)
##            except Exception:
##                client = chromadb.Client()
##
##            # Generate embedding (Ollama or fallback)
##            embedding = None
##            emb_source = None
##            try:
##                if ollama is None:
##                    raise RuntimeError('ollama package not available')
####                emb_resp = ollama.embeddings(model="mxbai-embed-large", prompt=file_text)
##                emb_resp = ollama.embeddings(model="mxbai-embed-large")
##                embedding = emb_resp.get('embedding')
##                emb_source = 'ollama'
##                if embedding is None:
##                    raise RuntimeError('Ollama returned no embedding')
##            except Exception:
##                from sentence_transformers import SentenceTransformer
##                st_model = SentenceTransformer('all-MiniLM-L6-v2')
##                vect = st_model.encode(file_text)
##                embedding = vect.tolist() if hasattr(vect, "tolist") else list(vect)
##                emb_source = 'sentence-transformers'
##
##            # Store in Chroma
##            doc_id = f"{user}_{safe_name}_{int(time.time())}"
##            collection = client.get_or_create_collection(name="docs")
##            collection.add(
##                ids=[doc_id],
##                embeddings=[embedding],
##                documents=[file_text]
##            )
##            print(f"collections is now before : (collection)")
##            print(f"Embeddings is now before : (embeddings)")
##
##            
##            try:
##                client.persist()
##            except Exception:
##                pass
##
##
##            # Log user query
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": f"[file_upload] {filename}",
##                "response": "[stored in RAG]",
##                "models": {"embed_model": emb_source}
##            }
##            log_user_query(user, entry)
##
##            # Notify frontend success
##            socketio.emit('state_update', {
##                'type': 'file_info',
##                'payload': {
##                    'message': f"Text '{filename}' processed for RAG",
##                    'doc_id': doc_id,
##                    'length': len(file_text),
##                    'rag_available': True,
##                    'embed_source': emb_source
##                },
##                'username': user
##            }, room=user)
##
##        except Exception as e:
##            socketio.emit('state_update', {
##                'type': 'file_error',
##                'payload': f'RAG processing failed: {e}',
##                'username': user
##            }, room=user)
##
##        return
##
##
##    # --- QUERY ---
##    if t == 'query':
##        text = payload if isinstance(payload, str) else payload.get('query', '')
##        print(f"[EVENT] query from user='{user}' text='{text}'")
##
##        if text:
##            text = text.replace(":","")
##            print(f"New TEXT query supplied : {text}")
##        else:
##            print("NO TEXT query supplied")
##            
##        result = ""
##
##        if any(x in text.lower() for x in ["this file", "this document", "this story", "this tale"]) and my_file_content:
##            try:
##                from rag_embeddings import file_embeddings
##
##                # initialize and build once per file
##                fe = file_embeddings()
##                fe.build_store(my_file_content)
##
##                # retrieve context chunks (already cleaned + action-line extracted in rag_embeddings.py)
##                context_chunks = fe.retrieve(query=text, top_k=3)
##
##                # Format them nicely for the model
##                if context_chunks:
##                    context = "Here are relevant action items from the document:\n"
##                    for chunk in context_chunks:
##                        context += f"- {chunk.strip()}\n"
##                else:
##                    context = "No relevant action items found in the document."
##
##                # now call Ollama chat with cleaned, formatted context
##                final_prompt = (
##                    f"User query: {text}\n\n"
##                    f"{context}\n\n"
##                    f"Answer:"
##                )
##                print(f"[QUERY] final_prompt is now : {final_prompt}")
##
##            except Exception as e:
##                print(f"[RAG ERROR QUERY] {e}")
##                result = ""
##
##            # Log and emit
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": final_prompt,
##                "response": None,
##                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
##            }
##            log_user_query(user, entry)
##            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
##
##            return
##
##        else:
##
##
##            # Log and emit
##            entry = {
##                "username": user,
##                "timestamp": datetime.now().isoformat(),
##                "query": text,
##                "response": None,
##                "models": payload.get('models', {}) if isinstance(payload, dict) else {}
##            }
##            log_user_query(user, entry)
##            socketio.emit('state_update', {'type': 'query', 'payload': entry, 'username': user})
##
##            return
##
##
### HTTP endpoints
##@app.route('/main_responses/<username>', methods=['GET'])
##def main_responses(username):
##    path = os.path.join(TMP_DIR, f"{username}_main.json")
##    if not os.path.exists(path):
##        return jsonify([])
##    with open(path, 'r', encoding='utf-8') as f:
##        try:
##            return jsonify(json.load(f))
##        except json.JSONDecodeError:
##            return jsonify([])
##
##@app.route('/user_logs/<username>', methods=["GET"])
##def get_user_logs(username):
##    logs = read_logs(username)
##    print(f"[DEBUG] Loaded logs for {username}: {len(logs)} entries")
##    return jsonify(logs)
##
##@app.post("/upload_audio")
##def upload_audio():
##    os.makedirs(TMP_DIR, exist_ok=True)
##    webm_path = os.path.join(TMP_DIR, "temp_audio.webm")
##    wav_path = os.path.join(TMP_DIR, "temp_audio.wav")
##
##    if 'audio' not in request.files:
##        return jsonify({"error": "No audio file"}), 400
##    request.files["audio"].save(webm_path)
##
##    # convert with ffmpeg if available
##    if os.system(f'ffmpeg -y -i "{webm_path}" -ar 16000 -ac 1 -f wav "{wav_path}"') != 0:
##        return jsonify({"error": "ffmpeg failed"}), 500
##
##    try:
##        text = whisper_model.transcribe(wav_path).get("text", "").strip()
##        cleaned = webui_listener.listen_text(text)
##    except Exception:
##        cleaned = ""
##
##    for p in (webm_path, wav_path):
##        try:
##            os.remove(p)
##        except:
##            pass
##
##    return jsonify({"transcript": cleaned})
##
##@app.post('/upload_file')
##def upload_file_http():
##    # HTTP multipart upload alternative to socket-based upload
##    if 'file' not in request.files:
##        return jsonify({"error": "No file provided"}), 400
##    file = request.files['file']
##    username = request.form.get('username', 'Tjaart')
##    filename = secure_filename(file.filename)
##    tmp_path = os.path.join(TMP_DIR, filename)
##    file.save(tmp_path)
##
##    # Try to detect type and if text, process into Chroma
##    mime_type, _ = mimetypes.guess_type(filename)
##
##    if mime_type and mime_type.startswith('image'):
##        print(f"[HTTP FILE] Image uploaded: {filename} by {username}")
##        socketio.emit('state_update', {'type': 'file_info', 'payload': f"Image '{filename}' received", 'username': username}, room=username)
##        return jsonify({"status": "ok", "type": "image"})
##
##    if mime_type and (mime_type.startswith('text') or filename.lower().endswith(('.txt', '.md', '.py'))):
##        print(f"[HTTP FILE] Text uploaded: {filename} by {username}")
##        try:
##            with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
##                file_text = f.read()
##        except Exception as e:
##            return jsonify({"error": f"Could not read file: {e}"}), 500
##
##        # Same RAG logic as socket handler
##        try:
##            client = None
##            if chromadb is not None:
##                try:
##                    from chromadb.config import Settings as ChromaSettings
##                    chroma_settings = ChromaSettings(
##                        chroma_db_impl="duckdb+parquet",
##                        persist_directory=os.path.join(BASE_DIR, "chroma_db")
##                    )
##                    client = chromadb.Client(chroma_settings)
##                except Exception:
##                    try:
##                        client = chromadb.Client()
##                    except Exception:
##                        client = None
##
##            embedding = None
##            emb_source = None
##            try:
##                if ollama is not None:
####                    emb_resp = ollama.embeddings(model="mxbai-embed-large", prompt=file_text)
##                    emb_resp = ollama.embeddings(model="mxbai-embed-large")
##                    embedding = emb_resp.get('embedding')
##                    emb_source = 'ollama'
##            except Exception:
##                try:
##                    from sentence_transformers import SentenceTransformer
##                    st_model = SentenceTransformer('all-MiniLM-L6-v2')
##                    vect = st_model.encode(file_text)
##                    embedding = vect.tolist() if hasattr(vect, 'tolist') else list(vect)
##                    emb_source = 'sentence-transformers'
##                except Exception:
##                    embedding = None
##                    emb_source = None
##
##            doc_id = f"{username}_{filename}_{int(time.time())}"
##            if embedding is not None and client is not None:
##                try:
##                    collection = client.get_or_create_collection(name="docs")
##                    collection.add(ids=[doc_id], embeddings=[embedding], documents=[file_text])
##                    try:
##                        client.persist()
##                    except Exception:
##                        pass
##                except Exception:
##                    client = None
##
##            if embedding is None or client is None:
##                fallback_index = os.path.join(TMP_DIR, "local_rag_index.json")
##                try:
##                    if os.path.exists(fallback_index):
##                        with open(fallback_index, 'r', encoding='utf-8') as fh:
##                            idx = json.load(fh)
##                    else:
##                        idx = []
##                except Exception:
##                    idx = []
##                idx_entry = {
##                    'doc_id': doc_id,
##                    'username': username,
##                    'filename': filename,
##                    'path': tmp_path,
##                    'length': len(file_text),
##                    'embedding_source': emb_source,
##                    'timestamp': datetime.now().isoformat()
##                }
##                idx.append(idx_entry)
##                try:
##                    with open(fallback_index, 'w', encoding='utf-8') as fh:
##                        json.dump(idx, fh, ensure_ascii=False, indent=2)
##                except Exception:
##                    pass
##
##            log_user_query(username, {"username": username, "timestamp": datetime.now().isoformat(), "query": f"[file_upload] {filename}", "response": "[stored in RAG]" if embedding is not None else "[saved - RAG unavailable]", "models": {"embed_model": emb_source or "none"}})
##
##    ##        prompt=f"Using this data: {data}. Respond to this prompt: {prompt}",
##            full_emb_prompt = f"Using this data: {response_text}. Respond to this prompt: {text}"
##            socketio.emit('state_update', {
##                'type': 'query',
##    ##            'payload': {'query': text, 'response': response_text},
##                'payload': full_emb_prompt,
##                'username': user
##            })
##
####            socketio.emit('state_update', {'type': 'file_info', 'payload': {'message': f"Text '{filename}' processed for RAG" if embedding is not None else f"Text '{filename}' saved (RAG unavailable)", 'doc_id': doc_id}, 'username': username}, room=username)
##            return jsonify({"status": "ok", "type": "text", "doc_id": doc_id})
##        except Exception as e:
##            print(f"[ERROR] HTTP RAG failed: {e}")
##            return jsonify({"error": str(e)}), 500
##
##    return jsonify({"status": "ok", "type": "unknown"})
##
##
##@app.route('/edge_tts', methods=['POST'])
##def edge_tts_endpoint():
##    """
##    Synthesize text with edge-tts and return MP3 bytes.
##    Frontend posts JSON: { text, voice?, style? }.
##    Returns: audio/mpeg binary (200) or JSON error (non-200).
##    """
##    data = request.get_json(silent=True) or {}
##    text = data.get('text') or data.get('response') or ""
##    voice = data.get('voice', 'en-US-GuyNeural')
##    style = data.get('style', None)  # currently unused in this simple example
##
##    if not text:
##        return jsonify({"error": "No 'text' provided"}), 400
##
##    # Try importing edge-tts
##    try:
##        import edge_tts
##    except Exception as e:
##        return jsonify({"error": f"edge-tts not installed or import failed: {e}"}), 500
##
##    # create a temporary file for the mp3
##    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
##    tmp_path = tmp.name
##    tmp.close()
##
##    async def synth_to_file(txt, voice_name):
##        # Use simple text; if you want SSML or style-based SSML, build it here.
##        comm = edge_tts.Communicate(txt, voice_name)
##        # save will write mp3 to path
##        await comm.save(tmp_path)
##
##    try:
##        # Run the async synth (blocks until complete)
##        asyncio.run(synth_to_file(text, voice))
##    except Exception as e:
##        # cleanup
##        try:
##            os.unlink(tmp_path)
##        except Exception:
##            pass
##        return jsonify({"error": f"edge-tts synthesis failed: {e}"}), 500
##
##    # Return the mp3 file bytes
##    try:
##        # send_file will set appropriate headers
##        resp = send_file(tmp_path, mimetype='audio/mpeg', as_attachment=False)
##        # Allow cross-origin if needed (your SocketIO allowed * already)
##        resp.headers['Access-Control-Allow-Origin'] = '*'
##        return resp
##    finally:
##        # cleanup file after response has been scheduled (Flask may still be reading it,
##        # but removing the file here is okay on many OSes; if you prefer, delete in a worker)
##        try:
##            os.unlink(tmp_path)
##        except Exception:
##            pass
##
##@app.route('/login', methods=["POST"])
##def login_user():
##    data = request.json
##    success, msg = do_login(data["username"], data["password"])
##    return jsonify({"success": success, "message": msg})
##
##@app.route('/signup', methods=["POST"])
##def signup_user():
##    data = request.json
##    success, msg = do_signup(data["username"], data["password"])
##    return jsonify({"success": success, "message": msg})
##
##@app.route('/users', methods=["GET"]) 
##def list_users():
##    return jsonify(get_users())
##
##@app.route('/delete_user', methods=["POST"]) 
##def remove_user():
##    data = request.json
##    success = delete_user(data["username"])
##    return jsonify({"success": success})
##
##@app.route('/', defaults={'path': ''})
##@app.route('/<path:path>')
##def serve(path):
##    if path != "" and os.path.exists(os.path.join(app.static_folder, path)):
##        return send_from_directory(app.static_folder, path)
##    return send_from_directory(app.static_folder, 'index.html')
##
##if __name__ == '__main__':
##    context = ('cert.pem', 'key.pem') if os.path.exists('cert.pem') and os.path.exists('key.pem') else None
##    socketio.run(app, host='0.0.0.0', port=5001, ssl_context=context)
##
